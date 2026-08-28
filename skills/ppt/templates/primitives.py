# -*- coding: utf-8 -*-
"""
ppt primitives — python-pptx 代码画 PPT 的积木箱
用法：复制本文件改造，或 import 后组装。数据区放项目事实（Step 1 取材结果）。

坐标约定：16:9（13.333 × 7.5 in）· 边距 0.55 · 页码 (12.35, 7.02)
"""
import math
import re

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── 色板（5+1 封闭系统；有项目品牌色则整体替换）────────────────────────
CORAL = "CC785C"       # 主色：强调条 / 关键数字 / 反色块
CORAL_DEEP = "9C4F37"  # 主色深：标题强调
CREAM = "FAF6F0"       # 底色：卡片斑马纹 / 次级底
PAPER = "FFFFFF"       # 纸白：卡片默认底
INK = "1C1917"         # 墨：正文标题
INK_SOFT = "57534E"    # 墨浅：正文
MUTED = "78716C"       # 灰：注释 / 页码（stone-500，白底 4.8:1 达 AA；暗底换 PAPER）
LINE = "E7E0D8"        # 线：卡片描边 / 分隔线

FONT_CN = "Microsoft YaHei"
FONT_MONO = "Consolas"
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = 0.55

C = RGBColor.from_string  # hex -> RGBColor

# ── 主题预设（均为对比度验证过的 5+1 封闭系统；deck 开始前 use_theme 切换）──
THEMES = {
    "warm":   dict(CORAL="CC785C", CORAL_DEEP="9C4F37", CREAM="FAF6F0", INK="1C1917",
                   INK_SOFT="57534E", MUTED="78716C", LINE="E7E0D8"),
    "tech":   dict(CORAL="3B6FE0", CORAL_DEEP="274CA0", CREAM="F4F6FA", INK="171B22",
                   INK_SOFT="4A5261", MUTED="6B7280", LINE="E2E6EE"),
    "forest": dict(CORAL="3E7C59", CORAL_DEEP="2C5B41", CREAM="F5F6F1", INK="1A1F1B",
                   INK_SOFT="4E574F", MUTED="6E776E", LINE="E3E6DE"),
}


def use_theme(name):
    """切换主题预设（改写模块级色常量；只影响之后的页面，deck 开始前调用）。
    warm 暖珊瑚（默认）/ tech 冷蓝 / forest 墨绿。有项目品牌色时优先直接改常量。"""
    for k, v in THEMES[name].items():
        globals()[k] = v

# ── 数据区（Step 1 取材结果，每项必须有出处）───────────────────────────
META = {
    "title": "PROJECT NAME",
    "subtitle": "项目说明 · Project Brief",
    "date": "YYYY-MM-DD",
    "author": "",
}


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def box(slide, x, y, w, h, fill=None, line=None):
    """矩形块：底色 + 描边（fill/line 是 hex 字符串，None = 透明）"""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = C(line); sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, content, size=13, color=INK, bold=False,
         align=PP_ALIGN.LEFT, font=FONT_CN, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """文本框。content: str | [str] | [(text, size, color, bold), ...]
    size 传 None 则用全局 size（tuple 项里各自覆盖）"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = content if isinstance(content, list) else [content]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(line, tuple):
            t, s, cl, b = line
        else:
            t, s, cl, b = line, size, color, bold
        r = p.add_run(); r.text = t
        r.font.size = Pt(s or size)
        r.font.color.rgb = C(cl)
        r.font.bold = b
        r.font.name = font
    return tb


def check_fit(content, size, w_in, h_in, spacing=1.0, label=""):
    """生成侧溢出预警：字符宽估算（中文≈1.0em / ASCII≈0.55em）× 雅黑行高 1.35em。
    保守估算而非精确排版，预警即拆行/拆页/砍字；最终以 Step 4 渲染核查为准。"""
    lines = content if isinstance(content, list) else [content]
    n = 0
    for ln in lines:
        t = ln if isinstance(ln, str) else ln[0]
        units = sum(1.0 if ord(ch) > 0x2E80 else 0.55 for ch in t)
        n += max(1, math.ceil(units * size / 72 / max(w_in - 0.2, 0.1)))  # 扣文本框左右 inset
    need = n * size * 1.35 * spacing / 72
    if need > h_in + 0.02:
        print(f"[FIT-WARN] {label or lines[0][:12]}: need {need:.2f}in > box {h_in:.2f}in ({n} lines)")
        return False
    return True


def page_chrome(slide, idx, label, badge=None):
    """每页统一 chrome：顶部 coral 短杠 + 章节 label + 右下页码 + 底部细线；
    badge = 右上角品牌角标（域名/项目名，Consolas 灰）"""
    box(slide, Inches(MARGIN), Inches(0.5), Inches(0.28), Inches(0.055), fill=CORAL)
    text(slide, Inches(MARGIN + 0.4), Inches(0.38), Inches(8), Inches(0.3),
         label.upper(), 10, MUTED, True)
    if badge:
        text(slide, Inches(10.0), Inches(0.38), Inches(2.78), Inches(0.3),
             badge, 10, MUTED, True, PP_ALIGN.RIGHT, font=FONT_MONO).name = "badge"
    text(slide, Inches(12.35), Inches(7.02), Inches(0.7), Inches(0.3),
         f"{idx:02d}", 10, MUTED, align=PP_ALIGN.RIGHT)
    box(slide, Inches(MARGIN), Inches(7.18), Inches(13.333 - 2 * MARGIN),
        Emu(9525), fill=LINE)


def shape_groups(slide, prefix):
    """按命名前缀取 shape 组：card0box/card0t/card0d 同属组 0（box 与其文字同进同退）。
    返回按组序排列的 [[shape...]...]，直接喂 animate.Anim.stagger(groups, 'float_up')。"""
    groups = {}
    for sp in slide.shapes:
        m = re.match(rf"^{re.escape(prefix)}(\d+)", sp.name)
        if m:
            groups.setdefault(int(m.group(1)), []).append(sp)
    return [groups[k] for k in sorted(groups)]


def notes(slide, script):
    """演讲者备注 = 口播稿：deck-to-video 与现场讲稿的数据源。每页 1~3 句，
    与页面同步写（Step 1 取材时同步产出），交付报告注明页数。"""
    slide.notes_slide.notes_text_frame.text = script


# ── 页面范式（组合优于发明；坐标用表达式算，不写魔法数）───────────────
def slide_cover(prs, meta):
    s = add_slide(prs)
    box(s, 0, 0, SW, SH, fill=CREAM)
    box(s, Inches(10.2), 0, Inches(3.133), SH, fill=CORAL)          # 右侧竖色块
    k = text(s, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4),
             meta.get("kicker", "PROJECT BRIEF").upper(), 13, CORAL, True)
    k.name = "kicker"
    t = text(s, Inches(0.9), Inches(2.15), Inches(8.6), Inches(2.2),
             [(meta["title"], 54, INK, True), (meta["subtitle"], 17, INK_SOFT, False)], spacing=1.05)
    t.name = "ctitle"
    d = box(s, Inches(0.95), Inches(4.45), Inches(2.2), Emu(19050), fill=CORAL)  # 分隔粗线
    d.name = "cdivider"
    m = text(s, Inches(0.9), Inches(4.85), Inches(8.4), Inches(1.6),
             [(meta.get("tagline", ""), 13, INK_SOFT, False),
              (f'{meta.get("domain", "")}    {meta.get("version", "")}    {meta["date"]}    {meta["author"]}',
               11, MUTED, False)], spacing=1.4)
    m.name = "cmeta"
    return s


def slide_cards(prs, idx, label, title, cards, cols=4, card_h=3.6):
    """1×cols 卡片页：cards = [(题, '描述行1\\n描述行2'), ...]"""
    s = add_slide(prs); page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    gap = 0.2
    card_w = (13.333 - 2 * MARGIN - gap * (cols - 1)) / cols
    for i, (t, d) in enumerate(cards):
        x = Inches(MARGIN + i * (card_w + gap))
        box(s, x, Inches(2.55), Inches(card_w), Inches(card_h), fill=PAPER, line=LINE).name = f"card{i}box"
        box(s, x, Inches(2.55), Inches(card_w), Inches(0.07), fill=CORAL).name = f"card{i}bar"  # 顶 coral 条
        text(s, x + Inches(0.25), Inches(2.9), Inches(card_w - 0.5), Inches(0.5),
             t, 19, INK, True).name = f"card{i}t"
        text(s, x + Inches(0.25), Inches(3.55), Inches(card_w - 0.4), Inches(2.4),
             d.split("\n"), 11.5, INK_SOFT, spacing=1.5).name = f"card{i}d"
    return s


def slide_numbers(prs, idx, label, title, nums, cols=3):
    """大数字墙：nums = [(数字, 标签, 注释), ...]，2 行 × cols"""
    s = add_slide(prs); page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    rows = (len(nums) + cols - 1) // cols
    gap = 0.28
    card_w = (13.333 - 2 * MARGIN - gap * (cols - 1)) / cols
    card_h = 2.0
    for i, (n, t, d) in enumerate(nums):
        row, col = divmod(i, cols)
        x = Inches(MARGIN + col * (card_w + gap))
        y = Inches(2.1 + row * (card_h + 0.3))
        box(s, x, y, Inches(card_w), Inches(card_h), fill=PAPER, line=LINE).name = f"num{i}box"
        text(s, x + Inches(0.28), y + Inches(0.22), Inches(card_w - 0.6), Inches(0.9),
             n, 40, CORAL_DEEP, True).name = f"num{i}v"
        text(s, x + Inches(0.3), y + Inches(1.18), Inches(card_w - 0.6), Inches(0.35),
             t, 13, INK, True).name = f"num{i}t"
        text(s, x + Inches(0.3), y + Inches(1.55), Inches(card_w - 0.55), Inches(0.35),
             d, 9.5, MUTED, font=FONT_MONO).name = f"num{i}d"
    return s


def slide_chain(prs, idx, label, title, nodes, hi=2, sub_title=None, subs=None):
    """横向链路图：nodes = [(名称, 子注释), ...]，hi = 反色强调的节点下标"""
    s = add_slide(prs); page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    gap = 0.37
    bw = (13.333 - 2 * MARGIN - gap * (len(nodes) - 1)) / len(nodes)
    for i, (t, d) in enumerate(nodes):
        x = Inches(MARGIN + i * (bw + gap))
        hot = i == hi
        box(s, x, Inches(2.3), Inches(bw), Inches(1.35),
            fill=CORAL if hot else PAPER, line=None if hot else LINE).name = f"node{i}box"
        text(s, x + Inches(0.12), Inches(2.52), Inches(bw - 0.24), Inches(0.4), t,
             12, PAPER if hot else INK, True, PP_ALIGN.CENTER).name = f"node{i}t"
        text(s, x + Inches(0.12), Inches(3.05), Inches(bw - 0.24), Inches(0.35), d,
             9, PAPER if hot else MUTED, hot, align=PP_ALIGN.CENTER,
             font=FONT_MONO).name = f"node{i}d"
        if i < len(nodes) - 1:
            text(s, x + Inches(bw - 0.02), Inches(2.62), Inches(gap + 0.06), Inches(0.5),
                 "→", 20, CORAL, True, PP_ALIGN.CENTER)
    if sub_title and subs:
        text(s, Inches(MARGIN), Inches(4.1), Inches(4), Inches(0.4), sub_title, 13, INK, True)
        sgap = 0.2
        scw = (13.333 - 2 * MARGIN - sgap * (len(subs) - 1)) / len(subs)
        for i, (u, t) in enumerate(subs):
            x = Inches(MARGIN + i * (scw + sgap))
            box(s, x, Inches(4.6), Inches(scw), Inches(1.15), fill=CREAM,
                line=LINE).name = f"sub{i}box"
            text(s, x + Inches(0.15), Inches(4.78), Inches(scw - 0.3), Inches(0.35),
                 t, 11.5, INK, True).name = f"sub{i}t"
            text(s, x + Inches(0.15), Inches(5.18), Inches(scw - 0.25), Inches(0.35),
                 u, 9, MUTED, font=FONT_MONO).name = f"sub{i}d"
    return s


def slide_rows(prs, idx, label, title, rows, col_split=None):
    """双栏清单（技术栈等）：rows = [(左标签, 左值, 右标签, 右值), ...]
    或单栏：rows = [(标签, 值), ...] + col_split=None 时每行全宽"""
    s = add_slide(prs); page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    if col_split:
        half = (len(rows) + 1) // 2
        for ri, (l1, v1, l2, v2) in enumerate(rows[:half]):
            y = Inches(1.95 + ri * 0.82)
            for ci, (lab, val) in enumerate([(l1, v1), (l2, v2)]):
                if not lab:
                    continue
                x = Inches(MARGIN + ci * 6.3)
                box(s, x, y, Inches(5.9), Inches(0.7),
                    fill=PAPER if ri % 2 == 0 else CREAM,
                    line=LINE).name = f"row{ri}c{ci}box"
                text(s, x + Inches(0.18), y + Inches(0.06), Inches(1.5), Inches(0.3),
                     lab, 9, MUTED).name = f"row{ri}c{ci}l"
                text(s, x + Inches(0.18), y + Inches(0.32), Inches(5.5), Inches(0.32),
                     val, 11.5, INK, True, font=FONT_MONO).name = f"row{ri}c{ci}v"
    else:
        for ri, (lab, val) in enumerate(rows):
            y = Inches(1.95 + ri * 0.82)
            box(s, Inches(MARGIN), y, Inches(12.23), Inches(0.7),
                fill=PAPER if ri % 2 == 0 else CREAM, line=LINE).name = f"row{ri}box"
            text(s, Inches(MARGIN + 0.18), y + Inches(0.06), Inches(1.8), Inches(0.3),
                 lab, 9, MUTED).name = f"row{ri}l"
            text(s, Inches(MARGIN + 0.18), y + Inches(0.32), Inches(11.5), Inches(0.32),
                 val, 11.5, INK, True).name = f"row{ri}v"
    return s


def slide_section(prs, no, title, points=None, total=None, domain=None):
    """章节分隔页（杂志风）：ink 反色整页 + hairline 上下框架 + 巨大序号 +
    章节进度 + 条目细线 + 右下超大暗纹序号。points 接 list（竖排条目，每条前
    coral 细线段）；total / domain 给出时显示「no / total」进度与左下域名。
    auto 编排：序号 wipe up → 标题/条目 fade 级联；框架与暗纹属背板不动。"""
    s = add_slide(prs)
    box(s, 0, 0, SW, SH, fill=INK).name = "secbg"
    # 上下 hairline 框架 + 顶栏 label / 进度（背板级，不进动画）
    box(s, Inches(MARGIN), Inches(0.62), Inches(13.333 - 2 * MARGIN),
        Emu(9525), fill=INK_SOFT).name = "secframe1"
    box(s, Inches(MARGIN), Inches(6.82), Inches(13.333 - 2 * MARGIN),
        Emu(9525), fill=INK_SOFT).name = "secframe2"
    text(s, Inches(MARGIN), Inches(0.32), Inches(6), Inches(0.3),
         "CHAPTER", 10, CORAL, True, font=FONT_MONO).name = "seclabel"
    if total:
        text(s, Inches(9.5), Inches(0.32), Inches(3.28), Inches(0.3),
             f"{no:02d} / {total:02d}", 10, PAPER, True,
             PP_ALIGN.RIGHT, font=FONT_MONO).name = "secprog"
    # 右下超大暗纹序号（负空间装饰，出血右缘）
    text(s, Inches(8.6), Inches(3.4), Inches(6.5), Inches(4.6),
         f"{no:02d}", 240, INK_SOFT, True, PP_ALIGN.RIGHT,
         font=FONT_MONO).name = "secghost"
    # 主内容：巨大序号 + 标题 + 条目
    text(s, Inches(0.82), Inches(1.35), Inches(5), Inches(2.2),
         f"{no:02d}", 110, CORAL, True, font=FONT_MONO).name = "secno"
    box(s, Inches(0.95), Inches(3.62), Inches(0.35), Emu(19050), fill=CORAL).name = "secbar"
    text(s, Inches(0.9), Inches(3.98), Inches(11), Inches(0.9),
         title, 40, PAPER, True).name = "sectitle"
    if points:
        pts = points if isinstance(points, list) else points.split("\n")
        for i, p in enumerate(pts):
            y = 5.05 + i * 0.46
            box(s, Inches(0.95), Inches(y + 0.12), Inches(0.22), Emu(9525),
                fill=CORAL).name = f"secdec{i}"
            text(s, Inches(1.35), Inches(y), Inches(10.5), Inches(0.36),
                 p, 13, CREAM).name = f"secpt{i}"
    text(s, Inches(MARGIN), Inches(7.0), Inches(8), Inches(0.3),
         domain or "", 9, PAPER, font=FONT_MONO).name = "secfoot"
    return s


def slide_timeline(prs, idx, label, title, events, hi=None, badge=None):
    """时间轴：events = [(时间点, 事件[换行最多两行]), ...]，横轴上下交替，
    hi = 高亮节点下标（当前位置）。auto 编排：节点按流向 wipe 0.15s 级联。"""
    s = add_slide(prs)
    page_chrome(s, idx, label, badge=badge)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    n = len(events)
    axis_y = 4.1
    box(s, Inches(MARGIN + 0.4), Inches(axis_y), Inches(12.23 - 0.8), Emu(9525),
        fill=LINE).name = "tlaxis"
    for i, (t, d) in enumerate(events):
        x = MARGIN + 0.7 + i * (12.23 - 1.4) / max(n - 1, 1)
        hot = i == hi
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.075), Inches(axis_y - 0.075),
                                 Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C(INK if hot else CORAL)
        dot.line.fill.background()
        dot.shadow.inherit = False
        dot.name = f"tl{i}dot"
        above = i % 2 == 0
        if above:
            text(s, Inches(x - 1.1), Inches(2.35), Inches(2.2), Inches(1.55),
                 [(d, 11.5, INK_SOFT, False), (t, 15, INK if hot else CORAL_DEEP, True)],
                 align=PP_ALIGN.CENTER, spacing=1.3,
                 anchor=MSO_ANCHOR.BOTTOM).name = f"tl{i}t"
        else:
            text(s, Inches(x - 1.1), Inches(axis_y + 0.25), Inches(2.2), Inches(1.55),
                 [(t, 15, INK if hot else CORAL_DEEP, True), (d, 11.5, INK_SOFT, False)],
                 align=PP_ALIGN.CENTER, spacing=1.3).name = f"tl{i}t"
    return s


def slide_versus(prs, idx, label, title, left, right, badge=None):
    """对比页：left/right = (标题, [行...])。左 PAPER 右 CREAM 双面板 + 中缝 VS 圆标
    （面板几何对称：两侧边距相等，VS 圆心正对页面中线）。auto 编排：左右两组 float_up 级联。"""
    s = add_slide(prs)
    page_chrome(s, idx, label, badge=badge)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    gap = 0.66
    pw = (12.23 - gap) / 2
    for ci, (head, rows, fill) in enumerate([(left[0], left[1], PAPER),
                                             (right[0], right[1], CREAM)]):
        x = Inches(MARGIN + ci * (pw + gap))
        box(s, x, Inches(2.1), Inches(pw), Inches(4.35), fill=fill,
            line=LINE).name = f"vs{ci}box"
        text(s, x + Inches(0.3), Inches(2.35), Inches(pw - 0.6), Inches(0.5),
             head, 19, INK, True).name = f"vs{ci}h"
        for ri, r_ in enumerate(rows):
            ry = Inches(3.1 + ri * 0.62)
            box(s, x + Inches(0.32), ry + Inches(0.1), Inches(0.12), Inches(0.12),
                fill=CORAL).name = f"vs{ci}r{ri}dot"
            text(s, x + Inches(0.6), ry, Inches(pw - 0.9), Inches(0.5),
                 r_, 11.5, INK_SOFT, spacing=1.25).name = f"vs{ci}r{ri}"
    vs = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.317), Inches(3.92), Inches(0.7), Inches(0.7))
    vs.fill.solid()
    vs.fill.fore_color.rgb = C(INK)
    vs.line.color.rgb = C(PAPER)
    vs.line.width = Pt(1.5)
    vs.shadow.inherit = False
    vs.name = "vsbadge"
    text(s, Inches(6.317), Inches(4.06), Inches(0.7), Inches(0.4), "VS", 15, PAPER, True,
         PP_ALIGN.CENTER).name = "vsbadge_t"
    return s


def slide_quote(prs, quote, source):
    """金句页：cream 整页，居中构图——coral 巨引号 → 大字引用 → 短杠 → mono 出处。
    auto 编排：引号 grow_turn → 正文/出处 fade。"""
    s = add_slide(prs)
    box(s, 0, 0, SW, SH, fill=CREAM).name = "qbg"
    text(s, Inches(6.167), Inches(1.35), Inches(1.0), Inches(1.5),
         '"', 100, CORAL, True, PP_ALIGN.CENTER, font=FONT_MONO).name = "qmark"
    text(s, Inches(1.67), Inches(3.05), Inches(10.0), Inches(1.6),
         quote, 30, INK, True, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         spacing=1.35).name = "qtext"
    box(s, Inches(6.517), Inches(4.95), Inches(0.3), Emu(19050), fill=CORAL).name = "qbar"
    text(s, Inches(1.67), Inches(5.3), Inches(10.0), Inches(0.4),
         source, 13, MUTED, True, PP_ALIGN.CENTER, font=FONT_MONO).name = "qsrc"
    return s


def slide_closing(prs, meta, footer_lines, slogan=None):
    """尾页：ink 反色整页 + coral logo 块。slogan 传 str 或
    [(文字, size, color, bold), ...]，缺省默认英文句 + domain。"""
    s = add_slide(prs)
    box(s, 0, 0, SW, SH, fill=INK)
    box(s, Inches(0.9), Inches(2.5), Inches(0.35), Emu(19050), fill=CORAL).name = "cdivider"
    if slogan is None:
        slogan_lines = [("Let's build something.", 40, PAPER, True),
                        (meta.get("domain", "asterforge.top"), 18, CORAL, False)]
    elif isinstance(slogan, str):
        slogan_lines = [(slogan, 40, PAPER, True),
                        (meta.get("domain", "asterforge.top"), 18, CORAL, False)]
    else:
        slogan_lines = slogan
    text(s, Inches(0.9), Inches(2.9), Inches(10), Inches(1.0),
         slogan_lines, spacing=1.25).name = "slogan"
    text(s, Inches(0.9), Inches(5.05), Inches(11), Inches(1.4), footer_lines,
         spacing=1.5, color=PAPER, font=FONT_MONO, size=12).name = "cfoot"
    box(s, Inches(11.9), Inches(0.9), Inches(1.3), Inches(1.3), fill=CORAL).name = "logobox"
    text(s, Inches(11.9), Inches(1.28), Inches(1.3), Inches(0.5), "TC", 24, PAPER, True,
         PP_ALIGN.CENTER).name = "logot"
    return s


def picture(slide, x, y, w, h, img, frame=True, caption=None):
    """品牌化图片：等比缩放居中入框（不拉伸），LINE 细边框，可选图注（Consolas 9.5 灰）。
    img 为本地路径（png/jpg）。返回 (picture, caption_tb|None)。"""
    pic = slide.shapes.add_picture(img, x, y)
    s = min(w / pic.width, h / pic.height)
    nw, nh = int(pic.width * s), int(pic.height * s)
    pic.left, pic.top = int(x + (w - nw) / 2), int(y + (h - nh) / 2)
    pic.width, pic.height = nw, nh
    if frame:
        pic.line.color.rgb = C(LINE)
        pic.line.width = Pt(0.75)
    pic.shadow.inherit = False
    pic.name = "media"
    cap = None
    if caption:
        cap = text(slide, x, y + h + 0.1, w, Inches(0.3), caption, 9.5, MUTED,
                   font=FONT_MONO)
        cap.name = "mediacap"
    return pic, cap


def slide_media(prs, idx, label, title, img, bullets, caption=None, img_side="left",
                sidebar=None, footnote=None):
    """图文页：图 + 要点列（img_side='right' 图换到右边）。
    bullets = [(要点, 说明), ...]，右侧逐条 float_up 级联（auto_page 可全代劳）。
    sidebar 模式（editorial）：图占六成 + 右窄栏圈注 [(短语, 说明), ...] + 底部脚注行，
    此时 bullets 忽略。"""
    s = add_slide(prs)
    page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    if sidebar:
        picture(s, Inches(MARGIN), Inches(1.95), Inches(7.4), Inches(4.25), img, caption=caption)
        for i, (t, d) in enumerate(sidebar):
            y = 2.0 + i * 1.08
            box(s, Inches(8.3), Inches(y + 0.08), Inches(0.14), Inches(0.14),
                fill=CORAL).name = f"media{i}dot"
            text(s, Inches(8.62), Inches(y), Inches(4.1), Inches(0.35),
                 t, 14, INK, True).name = f"media{i}h"
            check_fit(d, 10.5, 4.1, 0.62, label=f"media{i}d")
            text(s, Inches(8.62), Inches(y + 0.38), Inches(4.1), Inches(0.62),
                 d, 10.5, INK_SOFT, spacing=1.3).name = f"media{i}d"
        if footnote:
            text(s, Inches(MARGIN), Inches(6.55), Inches(12.2), Inches(0.3),
                 footnote, 9, MUTED, font=FONT_MONO).name = "mediafn"
        return s
    img_x = Inches(MARGIN) if img_side == "left" else Inches(6.9)
    lst_x = Inches(6.9) if img_side == "left" else Inches(MARGIN)
    picture(s, img_x, Inches(2.1), Inches(5.9), Inches(4.2), img, caption=caption)
    for i, (t, d) in enumerate(bullets):
        y = Inches(2.2 + i * 1.05)
        box(s, lst_x, y + Inches(0.09), Inches(0.14), Inches(0.14),
            fill=CORAL).name = f"media{i}dot"
        text(s, lst_x + Inches(0.35), y, Inches(5.4), Inches(0.4),
             t, 15, INK, True).name = f"media{i}h"
        text(s, lst_x + Inches(0.35), y + Inches(0.42), Inches(5.4), Inches(0.5),
             d, 11.5, INK_SOFT, spacing=1.3).name = f"media{i}d"
    return s


def slide_table(prs, idx, label, title, cols, rows, col_ws=None, row_h=0.52):
    """真表格页：cols = [列名...]，rows = [[c1, c2, ...], ...]，col_ws = 列宽比例列表
    （缺省等分；第一列建议窄做标签列）。列头 CREAM + 斑马行，6~9 行；tab* 命名吃
    auto 级联（行级快节奏）。单元格文本超列宽会被 check_fit 预警。"""
    s = add_slide(prs)
    page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    n = len(cols)
    ws = col_ws or [1 / n] * n
    tot = 13.333 - 2 * MARGIN
    xs, xe = [], MARGIN
    for w in ws:
        xs.append(xe)
        xe += tot * w
    box(s, Inches(MARGIN), Inches(1.95), Inches(tot), Inches(0.42),
        fill=CREAM).name = "tabhead"
    for j, c in enumerate(cols):
        text(s, Inches(xs[j] + 0.12), Inches(2.0), Inches(tot * ws[j] - 0.2), Inches(0.32),
             c, 9.5, MUTED, True, font=FONT_MONO).name = f"tabh{j}"
    for i, row in enumerate(rows):
        y = 2.37 + i * row_h
        box(s, Inches(MARGIN), Inches(y), Inches(tot), Inches(row_h),
            fill=PAPER if i % 2 == 0 else CREAM).name = f"tab{i}box"
        for j, cell in enumerate(row):
            check_fit(cell, 10.5, tot * ws[j] - 0.2, row_h - 0.08, label=f"tab{i}c{j}")
            text(s, Inches(xs[j] + 0.12), Inches(y + 0.04), Inches(tot * ws[j] - 0.2),
                 Inches(row_h - 0.08), cell, 10.5, INK if j == 0 else INK_SOFT, j == 0,
                 font=FONT_MONO if j == 0 else FONT_CN,
                 anchor=MSO_ANCHOR.MIDDLE).name = f"tab{i}c{j}"
    return s


def slide_dense(prs, idx, label, title, side, main_lines=None, img=None, img_h=3.1):
    """高密度分区页：主区约 60%（可选大图 + 大字陈述/大数字）+ 侧区 40% 分组细节清单。
    side = [(组名, [行...]), ...]；main_lines = [(text, size, color, bold), ...] 富文本。
    打破等分平铺：主区讲一件事，侧区把细节喂饱（每组 3~5 行 10.5pt）。"""
    s = add_slide(prs)
    page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
         title, 30, INK, True).name = "title"
    mw = 7.1
    y = 1.95
    if img:
        picture(s, Inches(MARGIN), Inches(y), Inches(mw), Inches(img_h), img)
        y += img_h + 0.25
    if main_lines:
        text(s, Inches(MARGIN), Inches(y), Inches(mw), Inches(6.8 - y),
             main_lines, spacing=1.35).name = "main0t"
    sx, sw = MARGIN + mw + 0.35, 13.333 - MARGIN - (MARGIN + mw + 0.35)
    yy = 1.95
    for gi, (h, lines) in enumerate(side):
        box(s, Inches(sx), Inches(yy + 0.07), Inches(0.16), Inches(0.16),
            fill=CORAL).name = f"side{gi}dot"
        text(s, Inches(sx + 0.3), Inches(yy), Inches(sw - 0.3), Inches(0.32),
             h, 12, CORAL_DEEP, True).name = f"side{gi}h"
        check_fit(lines, 10.5, sw - 0.3, 0.26 * len(lines), label=f"side{gi}d")
        text(s, Inches(sx + 0.3), Inches(yy + 0.38), Inches(sw - 0.3),
             Inches(0.26 * len(lines)), lines, 10.5, INK_SOFT, spacing=1.4).name = f"side{gi}d"
        yy += 0.38 + 0.26 * len(lines) + 0.28
    return s


def bar_chart(slide, x, y, w, h, cats, vals, title=None, horizontal=False, color=CORAL,
              highlight=None, highlight_color=INK):
    """可编辑原生柱/条图（单系列，非贴图）：cats 类目，vals 数值。
    horizontal=True 转横向条形图。highlight=i 时第 i 根强调为 highlight_color（其余 color）。
    品牌化：无 legend、数值标签 Consolas、值轴淡化。
    返回容器 GraphicFrame（供 anim.chart() 等按元素引用）。"""
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("s", vals)
    ct = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = slide.shapes.add_chart(ct, x, y, w, h, cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = title is not None
    if chart.has_title:
        chart.chart_title.text_frame.text = title
    plot = chart.plots[0]
    plot.gap_width = 60
    ser = plot.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = C(color)
    if highlight is not None:
        pt = ser.points[highlight]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = C(highlight_color)
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(11); dl.font.name = FONT_MONO; dl.font.color.rgb = C(INK)
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(12); ca.tick_labels.font.name = FONT_CN
    ca.format.line.color.rgb = C(LINE)
    va = chart.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = C(LINE)
    va.tick_labels.font.size = Pt(9); va.tick_labels.font.name = FONT_MONO
    va.tick_labels.font.color.rgb = C(MUTED)
    return gf


def line_chart(slide, x, y, w, h, cats, vals, title=None, color=CORAL):
    """趋势线（单系列原生可编辑，非贴图）：动画用 anim.chart(gf, 'series')——
    擦入即画线，符合时间序列直觉。末点数值标签 Consolas。返回 GraphicFrame。"""
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("s", vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, w, h, cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = title is not None
    if chart.has_title:
        chart.chart_title.text_frame.text = title
    ser = chart.plots[0].series[0]
    ser.format.line.color.rgb = C(color)
    ser.format.line.width = Pt(2.25)
    ser.smooth = False
    for pt_ in ser.points:  # 小实心标记，crisp 数据感
        pt_.marker.style = XL_MARKER_STYLE.CIRCLE
        pt_.marker.size = 6
        pt_.marker.format.fill.solid()
        pt_.marker.format.fill.fore_color.rgb = C(color)
        pt_.marker.format.line.color.rgb = C(color)
    dl = ser.points[len(vals) - 1].data_label  # 只标末点数值
    dl.text_frame.text = str(vals[-1])
    r = dl.text_frame.paragraphs[0].runs[0]
    r.font.size = Pt(14); r.font.bold = True; r.font.name = FONT_MONO
    r.font.color.rgb = C(INK)
    dl.position = XL_LABEL_POSITION.ABOVE
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(12); ca.tick_labels.font.name = FONT_CN
    ca.format.line.color.rgb = C(LINE)
    va = chart.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = C(LINE)
    va.tick_labels.font.size = Pt(9); va.tick_labels.font.name = FONT_MONO
    va.tick_labels.font.color.rgb = C(MUTED)
    gf.name = "chart"
    return gf


def donut_chart(slide, x, y, w, h, cats, vals, title=None, colors=None):
    """占比环（单系列原生可编辑）：逐块配色缺省 [CORAL, INK, MUTED, CORAL_DEEP, INK_SOFT]
    循环——只用封闭色板里的数据色（LINE/CREAM 是脚手架色，白底不可辨，禁入数据位）。
    右侧图例；中心留白可叠总数大字（用 text 叠加）。返回 GraphicFrame。"""
    if colors is None:
        colors = [CORAL, INK, MUTED, CORAL_DEEP, INK_SOFT]
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("s", vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, w, h, cd)
    chart = gf.chart
    chart.has_title = title is not None
    if chart.has_title:
        chart.chart_title.text_frame.text = title
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    chart.legend.font.name = FONT_CN
    ser = chart.plots[0].series[0]
    for i, pt_ in enumerate(ser.points):
        pt_.format.fill.solid()
        pt_.format.fill.fore_color.rgb = C(colors[i % len(colors)])
        pt_.format.line.color.rgb = C(PAPER)
        pt_.format.line.width = Pt(1.5)
    gf.name = "chart"
    return gf


def set_transition(prs, kind="fade"):
    """统一页面转场（XML 注入，save 前调用）：fade / push-left。已有转场的页跳过。"""
    import lxml.etree as etree

    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    body = '<p:push dir="l"/>' if kind == "push-left" else "<p:fade/>"
    xml = f'<p:transition xmlns:p="{P}" spd="med">{body}</p:transition>'
    for slide in prs.slides:
        el = slide._element
        if el.find(f"{{{P}}}transition") is None and el.find(f"{{{MC}}}AlternateContent") is None:
            # 幂等：已有转场（含 morph 的 AlternateContent 包装）的页跳过，防双 transition
            el.append(etree.fromstring(xml))


def narration_secs(slide, wps=4.2, pad=1.8, min_sec=3.0):
    """按口播稿估放映秒数：中文字数 / 每秒字数 + 停顿余量；无备注取 min_sec。"""
    t = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
    n = sum(1 for ch in t if not ch.isspace())
    return max(min_sec, n / wps + pad) if n else min_sec


def auto_show(prs, kind="fade", wps=4.2, pad=1.8, min_sec=3.0, override=None):
    """自动放映：每页按口播稿长度设转场 advTm 自动换页（save 前调用）。
    动画须为 after/with 链（auto_deck 默认即是）→ 整 deck 放着不管自动播完，录屏即视频草稿。
    override = {页号: 秒} 手动覆盖个别页；morph 页的 Choice/Fallback 转场都会写上。"""
    set_transition(prs, kind)
    PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    for i, s in enumerate(prs.slides, 1):
        sec = override.get(i) if override and i in override else narration_secs(
            s, wps, pad, min_sec)
        for el in s._element.iter(f"{{{PML}}}transition"):
            el.set("advTm", str(int(sec * 1000)))


P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
P159 = "http://schemas.microsoft.com/office/powerpoint/2015/09/main"
MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"


def _morph(slide, dur_ms=900):
    """给单页设 morph 平滑转场（PowerPoint 2019+；跨页同名 !!shape 补间）。
    morph 元素属 p159 扩展且须 mc:AlternateContent 包装——裸 p14:morph 会被 PowerPoint
    视为无效转场（EntryEffect 读作 0，带修改的保存直接丢弃）。"""
    import lxml.etree as etree

    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    for el in slide._element.findall(f"{{{P}}}transition"):
        slide._element.remove(el)  # 一页只允许一个 transition
    xml = (
        f'<mc:AlternateContent xmlns:mc="{MC}">'
        f'<mc:Choice xmlns:p159="{P159}" Requires="p159">'
        f'<p:transition xmlns:p="{P}" xmlns:p14="{P14}" spd="slow" p14:dur="{dur_ms}">'
        f'<p159:morph option="byObject"/></p:transition></mc:Choice>'
        f'<mc:Fallback><p:transition xmlns:p="{P}" spd="slow"><p:fade/></p:transition></mc:Fallback>'
        f'</mc:AlternateContent>'
    )
    slide._element.append(etree.fromstring(xml))


def growth_chart(prs, idx, label, title, cats, vals, color=CORAL, highlight=None,
                 horizontal=False, x=MARGIN, y=2.2, w=12.23, h=4.2, dur_ms=900):
    """morph 数据增长图：一次生成两页——零状态页（柱≈0/条≈0）+ 终态页。
    翻页时 PowerPoint 对 !! 同名元素做真补间：纵向柱子平滑长高、横向条形从左长出，
    数值随端点升起浮现。原生 chart 无逐柱补间（bldChart 只是擦入），数据增长叙事用本范式。
    horizontal=True 横向条形（类目左对齐列）。需 PowerPoint 2019+；占两页页码。
    返回 (零状态页, 终态页)。"""
    n = len(vals)
    gap = 0.25
    vmax = max(vals) or 1

    def build(page_idx, zero):
        s = add_slide(prs)
        page_chrome(s, page_idx, label)
        text(s, Inches(x), Inches(0.95), Inches(w), Inches(0.7), title, 30, INK, True)
        if not horizontal:
            bw = (w - gap * (n - 1)) / n
            base = box(s, Inches(x), Inches(y + h), Inches(w), Emu(9525), fill=LINE)
            base.name = "!!gbase"
            for i, (c, v) in enumerate(zip(cats, vals)):
                bx = x + i * (bw + gap)
                bh = 0.04 if zero else max(0.04, v / vmax * (h - 0.8))
                hot = i == highlight
                bar = box(s, Inches(bx), Inches(y + h - bh), Inches(bw), Inches(bh),
                          fill=LINE if zero else (INK if hot else color))
                bar.name = f"!!gbar{i}"
                lab = text(s, Inches(bx), Inches(y + h - bh - 0.42), Inches(bw), Inches(0.4),
                           str(v), 22, LINE if zero else (INK if hot else CORAL_DEEP), True,
                           PP_ALIGN.CENTER, font=FONT_MONO)
                lab.name = f"!!glab{i}"
                cat = text(s, Inches(bx), Inches(y + h + 0.08), Inches(bw), Inches(0.3),
                           c, 11.5, INK_SOFT, True, PP_ALIGN.CENTER)
                cat.name = f"!!gcat{i}"
        else:
            lab_w = 2.0
            area = w - lab_w - 0.9
            row_h = min(0.55, (h - 0.3 - gap * (n - 1)) / n)
            base = box(s, Inches(x + lab_w), Inches(y), Emu(9525), Inches(h), fill=LINE)
            base.name = "!!gbase"
            for i, (c, v) in enumerate(zip(cats, vals)):
                by = y + i * (row_h + gap)
                bw = 0.04 if zero else max(0.04, v / vmax * area)
                hot = i == highlight
                bar = box(s, Inches(x + lab_w), Inches(by), Inches(bw), Inches(row_h),
                          fill=LINE if zero else (INK if hot else color))
                bar.name = f"!!gbar{i}"
                lab = text(s, Inches(x + lab_w + bw + 0.08), Inches(by + row_h / 2 - 0.22),
                           Inches(0.9), Inches(0.44), str(v), 20,
                           LINE if zero else (INK if hot else CORAL_DEEP), True, font=FONT_MONO)
                lab.name = f"!!glab{i}"
                cat = text(s, Inches(x), Inches(by + row_h / 2 - 0.22), Inches(lab_w - 0.18),
                           Inches(0.44), c, 11.5, INK_SOFT, True, align=PP_ALIGN.RIGHT)
                cat.name = f"!!gcat{i}"
        return s

    s0 = build(idx, True)
    s1 = build(idx + 1, False)
    _morph(s1, dur_ms)
    return s0, s1


def growth_line(prs, idx, label, title, cats, vals, color=CORAL, highlight=None,
                x=MARGIN, y=2.1, w=12.23, h=3.9, dur_ms=900):
    """morph 趋势线增长：两页——零状态（折线贴基线）→ 终态（真实走势）。
    折线是等顶点 freeform，morph 对同构路径逐点插值 = 线从基线整体长出，
    数据点/数值标签随端点升起。原生 line_chart 的 bldChart 只是整图擦入，
    趋势叙事用本范式。需 PowerPoint 2019+；占两页页码。返回 (零状态页, 终态页)。"""
    n = len(vals)
    vmax = max(vals) or 1
    px = [x + i * (w / max(n - 1, 1)) for i in range(n)]

    def build(page_idx, zero):
        s = add_slide(prs)
        page_chrome(s, page_idx, label)
        text(s, Inches(x), Inches(0.95), Inches(w), Inches(0.7), title, 30, INK, True)
        box(s, Inches(x), Inches(y + h), Inches(w), Emu(9525), fill=LINE).name = "!!gaxis"
        ys = [y + h] * n if zero else [y + h - 0.3 - v / vmax * (h - 0.95) for v in vals]
        fb = s.shapes.build_freeform(Inches(px[0]), Inches(ys[0]), scale=1.0)
        fb.add_line_segments([(Inches(px[i]), Inches(ys[i])) for i in range(1, n)], close=False)
        ln = fb.convert_to_shape()
        ln.fill.background()
        ln.line.color.rgb = C(color)
        ln.line.width = Pt(2.5)
        ln.shadow.inherit = False
        ln.name = "!!gline"
        for i in range(n):
            hot = i == highlight
            d = 0.11
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px[i] - d / 2), Inches(ys[i] - d / 2),
                                     Inches(d), Inches(d))
            dot.fill.solid()
            dot.fill.fore_color.rgb = C(INK if hot else color)
            dot.line.fill.background()
            dot.shadow.inherit = False
            dot.name = f"!!gdot{i}"
            text(s, Inches(px[i] - 0.5), Inches(ys[i] - 0.55), Inches(1.0), Inches(0.4),
                 str(vals[i]), 12, LINE if zero else (INK if hot else CORAL_DEEP), True,
                 PP_ALIGN.CENTER, font=FONT_MONO).name = f"!!glab{i}"
            text(s, Inches(px[i] - 0.6), Inches(y + h + 0.12), Inches(1.2), Inches(0.3),
                 cats[i], 11.5, INK_SOFT, True, PP_ALIGN.CENTER).name = f"!!gcat{i}"
        return s

    s0 = build(idx, True)
    s1 = build(idx + 1, False)
    _morph(s1, dur_ms)
    return s0, s1


def growth_donut(prs, idx, label, title, cats, vals, colors=None, unit="",
                 x=MARGIN, y=1.9, w=12.23, h=4.5, dur_ms=900):
    """morph 占比环增长：两页——零状态（各扇区收缩成自己起始角上的细楔）→ 终态。
    扇区是等顶点 freeform 楔形，morph 逐点插值 = 各扇区从自己的起始边扫开、整环绽放。
    右侧图例 + 环心总数静态。原生 donut_chart 的 bldChart 只是整图擦入，构成叙事用
    本范式。需 PowerPoint 2019+；占两页页码。返回 (零状态页, 终态页)。"""
    if colors is None:
        colors = [CORAL, INK, MUTED, CORAL_DEEP, INK_SOFT]
    total = sum(vals)
    cx, cy = x + 2.9, y + h / 2
    R = min(h / 2 - 0.3, 2.3)
    rr = R * 0.62
    gap = math.radians(1.5)
    # 每扇区固定细分步数（零/终态同构，morph 才能逐点插值）
    segs, a = [], -math.pi / 2
    for v in vals:
        sweep = v / (total or 1) * 2 * math.pi
        segs.append((a + gap, a + sweep - gap, max(3, math.ceil(math.degrees(sweep) / 5))))
        a += sweep

    def wedge(s, a0, a1, steps, col):
        pts = [(cx + R * math.sin(a0 + (a1 - a0) * k / steps),
                cy - R * math.cos(a0 + (a1 - a0) * k / steps)) for k in range(steps + 1)]
        pts += [(cx + rr * math.sin(a1 - (a1 - a0) * k / steps),
                 cy - rr * math.cos(a1 - (a1 - a0) * k / steps)) for k in range(steps + 1)]
        fb = s.shapes.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]), scale=1.0)
        fb.add_line_segments([(Inches(px_), Inches(py_)) for px_, py_ in pts[1:]], close=True)
        sp = fb.convert_to_shape()
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(col)
        sp.line.color.rgb = C(PAPER)
        sp.line.width = Pt(1.5)
        sp.shadow.inherit = False
        return sp

    def build(page_idx, zero):
        s = add_slide(prs)
        page_chrome(s, page_idx, label)
        text(s, Inches(x), Inches(0.95), Inches(w), Inches(0.7), title, 30, INK, True)
        for i, (a0, a1, steps) in enumerate(segs):
            end = a0 + math.radians(0.5) if zero else a1
            wedge(s, a0, end, steps, colors[i % len(colors)]).name = f"!!gseg{i}"
        text(s, Inches(cx - 1.2), Inches(cy - 0.5), Inches(2.4), Inches(1.0),
             [(str(total) + unit, 34, INK, True), ("total", 11, MUTED, False)],
             align=PP_ALIGN.CENTER, spacing=1.05).name = "gtotal"
        for i, (c, v) in enumerate(zip(cats, vals)):
            ly = y + 0.6 + i * 0.78
            box(s, Inches(x + 6.3), Inches(ly), Inches(0.16), Inches(0.16),
                fill=colors[i % len(colors)]).name = f"gleg{i}dot"
            text(s, Inches(x + 6.65), Inches(ly - 0.08), Inches(3.6), Inches(0.35),
                 c, 13, INK, True).name = f"gleg{i}t"
            text(s, Inches(x + 10.3), Inches(ly - 0.08), Inches(1.5), Inches(0.35),
                 f"{v / (total or 1) * 100:.0f}%", 13, INK_SOFT, True, font=FONT_MONO
                 ).name = f"gleg{i}v"
        return s

    s0 = build(idx, True)
    s1 = build(idx + 1, False)
    _morph(s1, dur_ms)
    return s0, s1


# ── 读改现有 deck（改也是脚本改，禁手改 pptx；改完 save 新文件名） ──────
def _walk_shapes(shapes):
    for sp in shapes:
        if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(sp.shapes)
        else:
            yield sp


def deck_replace(prs, pairs, use_regex=False):
    """全文替换（runs 级，含备注页）：pairs = {旧: 新} 或 [(旧, 新), ...]。返回替换次数。"""
    import re as _re

    n = 0
    items = pairs.items() if isinstance(pairs, dict) else pairs
    for s in prs.slides:
        frames = [sp.text_frame for sp in _walk_shapes(s.shapes) if sp.has_text_frame]
        if s.has_notes_slide:
            frames.append(s.notes_slide.notes_text_frame)
        for tf in frames:
            for p in tf.paragraphs:
                for r in p.runs:
                    for old, new in items:
                        if use_regex:
                            r.text, k = _re.subn(old, new, r.text)
                        else:
                            k = r.text.count(old)
                            r.text = r.text.replace(old, new)
                        n += k
    return n


def deck_recolor(prs, mapping):
    """全 deck 换色（hex 大写字符串映射 {旧: 新}）：字体色 / 形状填充 / 线色。
    返回命中次数。配 THEMES 用：deck_recolor(prs, dict(zip(warm.values(), tech.values())))。"""
    n = 0

    def hit(color_obj):
        nonlocal n
        try:
            cur = str(color_obj.rgb)
        except Exception:
            return
        if cur in mapping:
            color_obj.rgb = C(mapping[cur])
            n += 1

    for s in prs.slides:
        for sp in _walk_shapes(s.shapes):
            if sp.has_text_frame:
                for p in sp.text_frame.paragraphs:
                    for r in p.runs:
                        hit(r.font.color)
            try:
                if sp.fill.type == MSO_FILL.SOLID:
                    hit(sp.fill.fore_color)
            except Exception:
                pass
            try:
                hit(sp.line.color)
            except Exception:
                pass
    return n


# ── 入口示例 ──────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_deck()
    slide_cover(prs, META)
    slide_cards(prs, 2, "Overview", "一个门面，四种角色", [
        ("展示", "5 个项目详情页\n章节式叙事"),
        ("接单", "小程序定制\n软硬件集成"),
        ("互动", "投票看板\n社区群组"),
        ("运营", "管理后台\n热更新"),
    ])
    slide_numbers(prs, 3, "Numbers", "用数字说话", [
        ("5", "详情页", "static pages"), ("25+", "API", "routes"), ("10", "模型", "models"),
        ("100", "A11y", "lighthouse"), ("60s", "热更新", "ISR"), ("0.008", "CLS", "web vitals"),
    ])
    slide_chain(prs, 4, "Architecture", "请求链路", [
        ("访客浏览器", "HTTPS 443"), ("nginx", "proxy"), ("Next.js", "standalone"),
        ("Prisma", "ORM"), ("SQLite", "file db"),
    ])
    s = add_slide(prs); page_chrome(s, 5, "Chart")
    bar_chart(s, Inches(0.55), Inches(1.6), Inches(8.0), Inches(4.8),
              ["Q1", "Q2", "Q3", "Q4"], [12, 25, 18, 31], title="decks per quarter")
    slide_closing(prs, META, ["github.com/example", "© 2026"])
    prs.save("deck-demo.pptx")
    print("OK deck-demo.pptx")
