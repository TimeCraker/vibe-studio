# -*- coding: utf-8 -*-
"""
animate.py — COM 动画后处理：按元素逐个声明进入动画（Windows + MS Office + pywin32）

用法（gen 脚本内，python-pptx 对象在手时）：
    from animate import Anim
    import primitives as P
    anim = Anim(prs)
    anim.fade(title, dur=0.5)                      # 淡入，0.5s 缓出
    anim.fx(card, 'float_up', dur=0.5, delay=0.12) # 通用入口：词表 10 种入场任选
    anim.stagger(P.shape_groups(s, 'card'), 'float_up', step=0.12)
                                                   # 级联：组内同进、组间 delay 递增
    anim.wipe(num_card, 'up')                      # 方向感：up / down / left / right
    anim.chart(chart_gf, by='category')            # 图表生长：allAtOnce / series / category
    prs.save('deck.pptx')
    anim.apply('deck.pptx')                        # COM 写入动画；之后才做 Step 4 渲染核查

节奏 > 花活：高级感来自时长（0.4~0.6s）+ 级联间隔（0.1~0.2s）+ 缓出，不是效果种类多。
效果词表全部 COM 实测（写入后 EffectType 读回一致）；90 年代花活（百叶窗/棋盘/螺旋）不收。
组件适配推荐见 SKILL.md Step 3 动画节。约定：apply 前必须 prs.save()（COM 按 shape id 定位）。
"""
import sys
from pathlib import Path

import pythoncom
import win32com.client

sys.stdout.reconfigure(errors="replace")  # GBK 控制台保险

FX = {  # MsoAnimEffect 子集（全部 COM 实测验证）
    "appear": 1, "fade": 10, "wipe": 22,      # 基础三件套
    "float_up": 39,    # Ascend，Float In 上浮——现代高级感主力
    "float_down": 42,  # Descend，Float In 下沉
    "zoom": 48,        # FadedZoom，柔缩放——大数字/强调块
    "grow_turn": 31,   # GrowAndTurn——图标/logo
    "ease_in": 29, "split": 16, "wheel": 21, "stretch": 17,  # 备选
}
DIR = {"up": 1, "down": 4, "right": 2, "left": 8}  # EffectParameters.Direction 方向值
TRIG = {"click": 0, "with": 2, "after": 3}         # MsoAnimTriggerType

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# 组件适配表（SKILL.md Step 3 的代码化）：(名字前缀, 效果, 方向, 时长, 级联间隔)
AUTO_RULES = [
    ("card", "float_up", None, 0.5, 0.12),
    ("num", "zoom", None, 0.6, 0.1),
    ("node", "wipe", "left", 0.4, 0.1),
    ("sub", "float_up", None, 0.5, 0.12),
    ("row", "float_up", None, 0.5, 0.12),
    ("media", "float_up", None, 0.5, 0.12),
    ("tl", "wipe", "left", 0.4, 0.15),
    ("vs", "float_up", None, 0.5, 0.15),
]
# 固定名字页面序列：(名字, 效果, 时长, delay, 方向|None)
COVER_SEQ = [("kicker", "fade", 0.5, 0.0, None), ("ctitle", "fade", 0.5, 0.25, None),
             ("cdivider", "fade", 0.5, 0.55, None), ("cmeta", "fade", 0.5, 0.7, None)]
CLOSING_SEQ = [("logobox", "grow_turn", 0.8, 0.0, None), ("logot", "grow_turn", 0.8, 0.0, None),
               ("slogan", "fade", 0.6, 0.2, None), ("cdivider", "fade", 0.5, 0.4, None),
               ("cfoot", "fade", 0.5, 0.5, None)]
SECTION_SEQ = [("secno", "wipe", 0.6, 0.0, "up"), ("sectitle", "fade", 0.5, 0.25, None),
               ("secbar", "fade", 0.4, 0.4, None), ("secpts", "fade", 0.4, 0.5, None)]
QUOTE_SEQ = [("qmark", "grow_turn", 0.6, 0.0, None), ("qtext", "fade", 0.6, 0.3, None),
             ("qbar", "fade", 0.4, 0.5, None), ("qsrc", "fade", 0.5, 0.6, None)]


class Anim:
    """收集 python-pptx shape 的动画声明，save 后一次性 COM 写入。"""

    def __init__(self, prs):
        self.prs = prs
        self._plan = []    # (页号, shape_id, effect_id, direction|None, trigger, dur, delay)
        self._charts = {}  # (页号, shape_id) -> by；shape_id 每页独立编号，必须带页做 key
        self._n = 0

    def _ref(self, sp):
        """定位 shape 的 (页号, shape_id)。按 XML 树身份找所属页——
        shape_id 每页独立编号不唯一，全 deck 按 id 搜会命中错误页的同 id 形状。"""
        el = sp._element
        while el is not None and not el.tag.endswith("}sld"):  # sp -> spTree -> cSld -> sld
            el = el.getparent()
        sld_el = el
        for i, s in enumerate(self.prs.slides, 1):
            if s._element is sld_el:
                return i, sp.shape_id
        raise ValueError(f"shape not in deck: id={sp.shape_id}")

    def _add(self, sp, name, direction, trigger, dur, delay):
        sn, sid = self._ref(sp)
        self._plan.append((sn, sid, FX[name], DIR.get(direction), TRIG[trigger], dur, delay))
        self._n += 1

    # ── 声明 API：动词方法覆盖高频，fx() 是通用入口 ──────────────────
    def fade(self, sp, trigger="after", dur=None, delay=None):
        self._add(sp, "fade", None, trigger, dur, delay)

    def wipe(self, sp, direction="up", trigger="after", dur=None, delay=None):
        self._add(sp, "wipe", direction, trigger, dur, delay)

    def appear(self, sp, trigger="after"):
        self._add(sp, "appear", None, trigger, None, None)

    def fx(self, sp, name, direction=None, trigger="after", dur=None, delay=None):
        """通用入口：name ∈ FX 词表（float_up / float_down / zoom / grow_turn /
        ease_in / split / wheel / stretch / fade / wipe / appear）。"""
        self._add(sp, name, direction, trigger, dur, delay)

    def stagger(self, groups, fx="float_up", step=0.12, dur=0.5, delay0=0.0, direction=None):
        """级联：groups = [[shape...]...]（P.shape_groups 取）。组内 with 同进、
        组间 delay 递增 step——真级联（重叠节奏），不是逐个等前一个放完。"""
        for gi, shapes in enumerate(groups):
            for si, sp in enumerate(shapes):
                trig = "after" if (gi == 0 and si == 0) else "with"
                self._add(sp, fx, direction, trig, dur, delay0 + gi * step)

    def chart(self, sp, by="category", trigger="after", dur=None, delay=None):
        """图表生长动画：by = allAtOnce / series / category。"""
        sn, sid = self._ref(sp)
        self._plan.append((sn, sid, FX["wipe"], DIR["up"], TRIG[trigger], dur, delay))
        self._charts[(sn, str(sid))] = by
        self._n += 1

    # ── auto 编排：组件适配表一键落地 ────────────────────────────────
    def auto_page(self, slide):
        """按组件适配表自动编排一页（primitives 系统命名是契约）：
        封面/尾页按固定序列级联；正文页 title fade -> 各命名组按推荐效果级联。
        未命名的 chrome/背板不动；morph 增长页（!!g* 命名）自然零命中。
        返回本页声明数；auto 后仍可手动 fx() 补充（如图表）。"""
        import primitives as P  # 同目录懒加载，取 shape_groups

        names = {sp.name for sp in slide.shapes}
        n0 = self._n
        seq = None
        if "kicker" in names:
            seq = COVER_SEQ
        elif "logobox" in names:
            seq = CLOSING_SEQ
        elif "secno" in names:
            seq = SECTION_SEQ
        elif "qmark" in names:
            seq = QUOTE_SEQ
        if seq:
            for i, (name, fxname, dur, delay, direction) in enumerate(seq):
                if name in names:
                    self.fx(next(sp for sp in slide.shapes if sp.name == name), fxname,
                            direction=direction,
                            trigger="after" if i == 0 else "with", dur=dur,
                            delay=delay if i else None)
        else:
            if "title" in names:
                self.fx(next(sp for sp in slide.shapes if sp.name == "title"),
                        "fade", dur=0.4)
            if "media" in names:  # 图文页主图：标题后淡入，再接要点级联
                self.fx(next(sp for sp in slide.shapes if sp.name == "media"),
                        "fade", dur=0.5)
            for prefix, fxname, direction, dur, step in AUTO_RULES:
                groups = P.shape_groups(slide, prefix)
                if groups:
                    self.stagger(groups, fxname, step=step, dur=dur, direction=direction)
        return self._n - n0

    def auto_deck(self):
        """全 deck 自动编排。打印每页声明摘要供核对；返回总声明数。"""
        total = 0
        parts = []
        for i, s in enumerate(self.prs.slides, 1):
            n = self.auto_page(s)
            total += n
            parts.append(f"p{i}:{n}")
        print(f"auto-deck: {' '.join(parts)} = {total} effects")
        return total

    def apply(self, path):
        """打开已保存的 pptx 写入全部动画与节奏；图表生长另做 XML 收尾；
        全部落盘后重开读回核对（数量/页位/类型），防 COM 静默丢失。"""
        p = str(Path(path).resolve())
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch("PowerPoint.Application")
        try:
            pres = app.Presentations.Open(p, False, False, False)
            for slide_no, spid, fxid, direction, trig, dur, delay in self._plan:
                sld = pres.Slides(slide_no)
                eff = sld.TimeLine.MainSequence.AddEffect(_by_id(sld, spid), fxid, 0, trig)
                if direction is not None:
                    eff.EffectParameters.Direction = direction
                if dur or delay:
                    t = eff.Timing
                    if dur:
                        t.Duration = dur
                        t.SmoothEnd = True  # 缓出——高级感的一半
                    if delay:
                        t.TriggerDelayTime = delay
            try:
                pres.Save()
            except Exception:
                import time
                time.sleep(1.5)   # 偶发 sharing violation（杀软/索引器短暂锁文件），重试一次
                pres.Save()
            pres.Close()
        finally:
            app.Quit()
        if self._charts:
            _retag_charts(p, self._charts)
        _verify_written(p, self._plan)
        print(f"OK animated {self._n} effects, readback verified -> {p}")


def _by_id(sld, spid):
    for i in range(1, sld.Shapes.Count + 1):
        sh = sld.Shapes(i)
        if sh.Id == spid:
            return sh
    raise KeyError(f"shape id {spid} not on slide")


def _verify_written(path, plan):
    """apply 收尾自检：重开 pptx 读回 MainSequence，核对每页效果数量与 (shape id, 类型)。
    不齐 = COM 静默丢弃过动画，立即报错让 gen 脚本停下，不带着假 OK 交付。"""
    import pythoncom
    import win32com.client

    want = {}
    for slide_no, spid, fxid, *_ in plan:
        want.setdefault(slide_no, []).append((spid, fxid))
    pythoncom.CoInitialize()
    app = win32com.client.Dispatch("PowerPoint.Application")
    problems = []
    try:
        pres = app.Presentations.Open(str(Path(path).resolve()), True, False, False)
        try:
            for sn, expect in want.items():
                seq = pres.Slides(sn).TimeLine.MainSequence
                got = [(seq(i).Shape.Id, seq(i).EffectType) for i in range(1, seq.Count + 1)]
                if sorted(got) != sorted(expect):
                    missing = sorted(set(expect) - set(got))
                    extra = sorted(set(got) - set(expect))
                    problems.append(f"slide {sn}: missing {missing[:3]} extra {extra[:3]}")
        finally:
            pres.Close()
    finally:
        app.Quit()
    if problems:
        raise RuntimeError("animation readback mismatch:\n" + "\n".join(problems))


def _retag_charts(path, charts):
    """COM 默认写 bldAsOne（整图一动）；替换为 bldChart 按类目/系列生长。
    charts 的 key 是 (页号, shape_id str)——spid 每页独立编号，跨页同名 id 不碰撞。"""
    import lxml.etree as etree

    from pptx import Presentation

    prs = Presentation(path)
    for si, slide in enumerate(prs.slides, 1):
        for bg in slide._element.iter(f"{{{P_NS}}}bldGraphic"):
            by = charts.get((si, bg.get("spid")))
            if not by:
                continue
            for one in bg.findall(f"{{{P_NS}}}bldAsOne"):
                bg.remove(one)
            sub = etree.SubElement(bg, f"{{{P_NS}}}bldSub")
            etree.SubElement(sub, f"{{{A_NS}}}bldChart").set("bld", by)
    prs.save(path)
