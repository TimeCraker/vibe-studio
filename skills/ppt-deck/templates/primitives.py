# -*- coding: utf-8 -*-
"""
ppt-deck primitives — python-pptx 代码画 PPT 的积木箱
用法：复制本文件改造，或 import 后组装。数据区放项目事实（Step 1 取材结果）。

坐标约定：16:9（13.333 × 7.5 in）· 边距 0.55 · 页码 (12.35, 7.02)
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── 色板（5+1 封闭系统；有项目品牌色则整体替换）────────────────────────
CORAL = "CC785C"       # 主色：强调条 / 关键数字 / 反色块
CORAL_DEEP = "9C4F37"  # 主色深：标题强调
CREAM = "FAF6F0"       # 底色：卡片斑马纹 / 次级底
PAPER = "FFFFFF"       # 纸白：卡片默认底
INK = "1C1917"         # 墨：正文标题
INK_SOFT = "57534E"    # 墨浅：正文
MUTED = "A8A29E"       # 灰：注释 / 页码（白底用，暗底换 PAPER）
LINE = "E7E0D8"        # 线：卡片描边 / 分隔线

FONT_CN = "Microsoft YaHei"
FONT_MONO = "Consolas"
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = 0.55

C = RGBColor.from_string  # hex -> RGBColor

# ── 数据区（Step 1 取材结果，每项必须有出处）───────────────────────────
META = {
    "title": "PROJECT NAME",
    "subtitle": "项目说明 · Project Brief",
    "date": "YYYY-MM-DD",
    "author": "",
}


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def box(slide, x, y, w, h, fill=None, line=None):
    """矩形块：底色 + 描边（fill/line 是 hex 字符串，None = 透明）"""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = C(line); sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, content, size=13, color=INK, bold=False,
         align=PP_ALIGN.LEFT, font=FONT_CN, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """文本框。content: str | [str] | [(text, size, color, bold), ...]
    size 传 None 则用全局 size（tuple 项里各自覆盖）"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = content if isinstance(content, list) else [content]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(line, tuple):
            t, s, cl, b = line
        else:
            t, s, cl, b = line, size, color, bold
        r = p.add_run(); r.text = t
        r.font.size = Pt(s or size)
        r.font.color.rgb = C(cl)
        r.font.bold = b
        r.font.name = font
    return tb


def page_chrome(slide, idx, label):
    """每页统一 chrome：顶部 coral 短杠 + 章节 label + 右下页码 + 底部细线"""
    box(slide, Inches(MARGIN), Inches(0.5), Inches(0.28), Inches(0.055), fill=CORAL)
    text(slide, Inches(MARGIN + 0.4), Inches(0.38), Inches(8), Inches(0.3),
         label.upper(), 10, MUTED, True)
    text(slide, Inches(12.35), Inches(7.02), Inches(0.7), Inches(0.3),
         f"{idx:02d}", 10, MUTED, align=PP_ALIGN.RIGHT)
    box(slide, Inches(MARGIN), Inches(7.18), Inches(13.333 - 2 * MARGIN),
        Emu(9525), fill=LINE)


# ── 页面范式（组合优于发明；坐标用表达式算，不写魔法数）───────────────
def slide_cover(prs, meta):
    s = add_slide(prs)
    box(s, 0, 0, SW, SH, fill=CREAM)
    box(s, Inches(10.2), 0, Inches(3.133), SH, fill=CORAL)          # 右侧竖色块
    box(s, Inches(9.85), Inches(2.1), Inches(0.35), Inches(0.07), fill=INK)
    text(s, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), "PROJECT BRIEF", 13, CORAL, True)
    text(s, Inches(0.9), Inches(2.15), Inches(8.6), Inches(2.2),
         [(meta["title"], 54, INK, True), (meta["subtitle"], 17, INK_SOFT, False)], spacing=1.05)
    box(s, Inches(0.95), Inches(4.45), Inches(2.2), Emu(19050), fill=CORAL)  # 分隔粗线
    text(s, Inches(0.9), Inches(4.85), Inches(8.4), Inches(1.6),
         [(meta.get("tagline", ""), 13, INK_SOFT, False),
          (f'{meta.get("domain", "")}    {meta.get("version", "")}    {meta["date"]}    {meta["author"]}',
           11, MUTED, False)], spacing=1.4)
    return s


def slide_cards(prs, idx, label, title, cards, cols=4, card_h=3.6):
    """1×cols 卡片页：cards = [(题, '描述行1\\n描述行2'), ...]"""
    s = add_slide(prs); page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7), title, 30, INK, True)
    gap = 0.2
    card_w = (13.333 - 2 * MARGIN - gap * (cols - 1)) / cols
    for i, (t, d) in enumerate(cards):
        x = Inches(MARGIN + i * (card_w + gap))
        box(s, x, Inches(2.55), Inches(card_w), Inches(card_h), fill=PAPER, line=LINE)
        box(s, x, Inches(2.55), Inches(card_w), Inches(0.07), fill=CORAL)     # 顶 coral 条
        text(s, x + Inches(0.25), Inches(2.9), Inches(card_w - 0.5), Inches(0.5), t, 19, INK, True)
        text(s, x + Inches(0.25), Inches(3.55), Inches(card_w - 0.4), Inches(2.4),
             d.split("\n"), 11.5, INK_SOFT, spacing=1.5)
    return s


def slide_numbers(prs, idx, label, title, nums, cols=3):
    """大数字墙：nums = [(数字, 标签, 注释), ...]，2 行 × cols"""
    s = add_slide(prs); page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7), title, 30, INK, True)
    rows = (len(nums) + cols - 1) // cols
    gap = 0.28
    card_w = (13.333 - 2 * MARGIN - gap * (cols - 1)) / cols
    card_h = 2.0
    for i, (n, t, d) in enumerate(nums):
        row, col = divmod(i, cols)
        x = Inches(MARGIN + col * (card_w + gap))
        y = Inches(2.1 + row * (card_h + 0.3))
        box(s, x, y, Inches(card_w), Inches(card_h), fill=PAPER, line=LINE)
        text(s, x + Inches(0.28), y + Inches(0.22), Inches(card_w - 0.6), Inches(0.9),
             n, 40, CORAL_DEEP, True)
        text(s, x + Inches(0.3), y + Inches(1.18), Inches(card_w - 0.6), Inches(0.35), t, 13, INK, True)
        text(s, x + Inches(0.3), y + Inches(1.55), Inches(card_w - 0.55), Inches(0.35),
             d, 9.5, MUTED, font=FONT_MONO)
    return s


def slide_chain(prs, idx, label, title, nodes, hi=2, sub_title=None, subs=None):
    """横向链路图：nodes = [(名称, 子注释), ...]，hi = 反色强调的节点下标"""
    s = add_slide(prs); page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7), title, 30, INK, True)
    gap = 0.37
    bw = (13.333 - 2 * MARGIN - gap * (len(nodes) - 1)) / len(nodes)
    for i, (t, d) in enumerate(nodes):
        x = Inches(MARGIN + i * (bw + gap))
        hot = i == hi
        box(s, x, Inches(2.3), Inches(bw), Inches(1.35),
            fill=CORAL if hot else PAPER, line=None if hot else LINE)
        text(s, x + Inches(0.12), Inches(2.52), Inches(bw - 0.24), Inches(0.4), t,
             12, PAPER if hot else INK, True, PP_ALIGN.CENTER)
        text(s, x + Inches(0.12), Inches(3.05), Inches(bw - 0.24), Inches(0.35), d,
             9, PAPER if hot else MUTED, align=PP_ALIGN.CENTER, font=FONT_MONO)
        if i < len(nodes) - 1:
            text(s, x + Inches(bw - 0.02), Inches(2.62), Inches(gap + 0.06), Inches(0.5),
                 "→", 20, CORAL, True, PP_ALIGN.CENTER)
    if sub_title and subs:
        text(s, Inches(MARGIN), Inches(4.1), Inches(4), Inches(0.4), sub_title, 13, INK, True)
        sgap = 0.2
        scw = (13.333 - 2 * MARGIN - sgap * (len(subs) - 1)) / len(subs)
        for i, (u, t) in enumerate(subs):
            x = Inches(MARGIN + i * (scw + sgap))
            box(s, x, Inches(4.6), Inches(scw), Inches(1.15), fill=CREAM, line=LINE)
            text(s, x + Inches(0.15), Inches(4.78), Inches(scw - 0.3), Inches(0.35), t, 11.5, INK, True)
            text(s, x + Inches(0.15), Inches(5.18), Inches(scw - 0.25), Inches(0.35),
                 u, 9, MUTED, font=FONT_MONO)
    return s


def slide_rows(prs, idx, label, title, rows, col_split=None):
    """双栏清单（技术栈等）：rows = [(左标签, 左值, 右标签, 右值), ...]
    或单栏：rows = [(标签, 值), ...] + col_split=None 时每行全宽"""
    s = add_slide(prs); page_chrome(s, idx, label)
    text(s, Inches(MARGIN), Inches(0.95), Inches(12.2), Inches(0.7), title, 30, INK, True)
    if col_split:
        half = (len(rows) + 1) // 2
        for ri, (l1, v1, l2, v2) in enumerate(rows[:half]):
            y = Inches(1.95 + ri * 0.82)
            for ci, (lab, val) in enumerate([(l1, v1), (l2, v2)]):
                if not lab:
                    continue
                x = Inches(MARGIN + ci * 6.3)
                box(s, x, y, Inches(5.9), Inches(0.7), fill=PAPER if ri % 2 == 0 else CREAM, line=LINE)
                text(s, x + Inches(0.18), y + Inches(0.06), Inches(1.5), Inches(0.3), lab, 9, MUTED)
                text(s, x + Inches(0.18), y + Inches(0.32), Inches(5.5), Inches(0.32),
                     val, 11.5, INK, True, font=FONT_MONO)
    else:
        for ri, (lab, val) in enumerate(rows):
            y = Inches(1.95 + ri * 0.82)
            box(s, Inches(MARGIN), y, Inches(12.23), Inches(0.7),
                fill=PAPER if ri % 2 == 0 else CREAM, line=LINE)
            text(s, Inches(MARGIN + 0.18), y + Inches(0.06), Inches(1.8), Inches(0.3), lab, 9, MUTED)
            text(s, Inches(MARGIN + 0.18), y + Inches(0.32), Inches(11.5), Inches(0.32),
                 val, 11.5, INK, True)
    return s


def slide_closing(prs, meta, footer_lines):
    """尾页：ink 反色整页 + coral logo 块"""
    s = add_slide(prs)
    box(s, 0, 0, SW, SH, fill=INK)
    box(s, Inches(0.9), Inches(2.35), Inches(0.35), Emu(19050), fill=CORAL)
    text(s, Inches(0.9), Inches(2.75), Inches(10), Inches(1.0),
         [("Let's build something.", 40, PAPER, True),
          (meta.get("domain", "asterforge.top"), 18, CORAL, False)], spacing=1.25)
    text(s, Inches(0.9), Inches(4.9), Inches(11), Inches(1.4), footer_lines,
         spacing=1.5, color=MUTED, font=FONT_MONO, size=12)
    box(s, Inches(11.9), Inches(0.9), Inches(1.3), Inches(1.3), fill=CORAL)
    text(s, Inches(11.9), Inches(1.28), Inches(1.3), Inches(0.5), "TC", 24, PAPER, True, PP_ALIGN.CENTER)
    return s


# ── 入口示例 ──────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_deck()
    slide_cover(prs, META)
    slide_cards(prs, 2, "Overview", "一个门面，四种角色", [
        ("展示", "5 个项目详情页\n章节式叙事"),
        ("接单", "小程序定制\n软硬件集成"),
        ("互动", "投票看板\n社区群组"),
        ("运营", "管理后台\n热更新"),
    ])
    slide_numbers(prs, 3, "Numbers", "用数字说话", [
        ("5", "详情页", "static pages"), ("25+", "API", "routes"), ("10", "模型", "models"),
        ("100", "A11y", "lighthouse"), ("60s", "热更新", "ISR"), ("0.008", "CLS", "web vitals"),
    ])
    slide_chain(prs, 4, "Architecture", "请求链路", [
        ("访客浏览器", "HTTPS 443"), ("nginx", "proxy"), ("Next.js", "standalone"),
        ("Prisma", "ORM"), ("SQLite", "file db"),
    ])
    slide_closing(prs, META, ["github.com/example", "© 2026"])
    prs.save("deck-demo.pptx")
    print("OK deck-demo.pptx")
