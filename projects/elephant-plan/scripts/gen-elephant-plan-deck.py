# -*- coding: utf-8 -*-
"""elephant-e-ai 2a 冲刺排期图（2 页）。

P1 依赖分层图：每张卡显式标注前置 + 依赖箭头（相邻层实线 / 跨层虚线 / 关键路径粗黑线）。
P2 盯三件事。
数据源：docs/openspec-sdd-plan.md §2 依赖表（origin/master 22c5ba4，2026-09-04）。
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from primitives import (  # noqa: E402
    C, CREAM, INK, INK_SOFT, LINE, MUTED, PAPER, PP_ALIGN, MSO_ANCHOR,
    Inches, Pt, add_slide, box, check_fit, new_deck, page_chrome, text,
    use_theme,
)
from pptx.enum.shapes import MSO_CONNECTOR  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

use_theme("tech")

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "elephant-2a-plan.pptx")
B_COLOR, A_COLOR = "3B6FE0", "9C4F37"   # B线=你(蓝) A线=队友(红) 灰=弹性件
CP_COLOR = "171B22"                      # 关键路径粗黑线（避开两线框色）

try:
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
except ImportError:  # 兼容旧版 python-pptx
    MSO_LINE_DASH_STYLE = None


def dep_card(s, x, y, w, h, title, prereq, border, fill=CREAM):
    box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill, line=border)
    text(s, Inches(x + 0.07), Inches(y + 0.04), Inches(w - 0.14), Inches(h - 0.08),
         [(title, 13, INK, True), (prereq, 9, INK_SOFT, False)],
         anchor=MSO_ANCHOR.MIDDLE, spacing=1.08)
    check_fit(title, 13, w - 0.14, h, label="title")


def arrow(s, x0, y0, x1, y1, color=None, width=1.0, dashed=False):
    """依赖箭头：起点=前置卡底边，终点=后继卡边框，带三角箭头。"""
    conn = s.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0), Inches(x1), Inches(y1))
    conn.line.color.rgb = C(color or INK_SOFT)
    conn.line.width = Pt(width)
    if dashed and MSO_LINE_DASH_STYLE is not None:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    conn.shadow.inherit = False
    return conn


prs = new_deck()

# ══ 第 1 页：依赖分层图（每卡标前置）═════════════════════════════════
s = add_slide(prs)
page_chrome(s, 1, "elephant-e-ai · 2a 冲刺排期")
text(s, Inches(0.55), Inches(0.62), Inches(7.9), Inches(0.5),
     "依赖分层图：每项开工前必须先有什么", size=30, color=INK, bold=True)
text(s, Inches(0.55), Inches(1.18), Inches(7.6), Inches(0.3),
     "自上而下 5 层 · 箭头起点=前置，终点=可开工项", size=11, color=MUTED)

# 图例：线样 + 说明（右上一行线样，下一行框色）
arrow(s, 8.35, 1.24, 8.75, 1.24)
text(s, Inches(8.80), Inches(1.13), Inches(0.75), Inches(0.24), "相邻层", size=9, color=MUTED)
arrow(s, 9.70, 1.24, 10.10, 1.24, dashed=True)
text(s, Inches(10.15), Inches(1.13), Inches(0.62), Inches(0.24), "跨层", size=9, color=MUTED)
arrow(s, 11.00, 1.24, 11.40, 1.24, color=CP_COLOR, width=2.25)
text(s, Inches(11.45), Inches(1.13), Inches(1.33), Inches(0.24), "关键路径", size=9, color=MUTED)
text(s, Inches(8.35), Inches(1.40), Inches(4.43), Inches(0.24),
     "蓝框 B线·你　红框 A线·队友　灰框 弹性件", size=9, color=MUTED)

W, H = 1.9, 0.72
LY = [1.80, 2.85, 3.90, 4.95, 6.00]     # 第0~4层卡片顶
CX = {"1-6": 2.45, "1-4": 4.5, "2-6": 8.6, "2-5": 10.65,
      "2-1": 4.5, "2-10": 2.45, "2-11": 4.5, "4-1": 6.55, "2-12": 8.6}

for i, lab in enumerate(["第0层", "第1层", "第2层", "第3层", "第4层"]):
    text(s, Inches(0.55), Inches(LY[i] + 0.10), Inches(1.0), Inches(0.3),
         [(lab, 9, MUTED, True)], spacing=1.1)

CARDS = [  # (层, key, 标题, 前置, 框色)
    (0, "1-6", "1-6 密钥托管", "前置：无 · 谁有空谁做", MUTED),
    (0, "1-4", "1-4 租户上下文", "前置：无 · 最小版先行", A_COLOR),
    (0, "2-6", "2-6 指令通道", "前置：无 · 今天开工", B_COLOR),
    (0, "2-5", "2-5 审计待办", "前置：无 · 4-20 前落地", MUTED),
    (1, "2-1", "2-1 主数据补齐", "前置：1-4", A_COLOR),
    (2, "2-10", "2-10 卡池", "前置：2-1 ＋ 1-6", A_COLOR),
    (2, "2-11", "2-11 花名册", "前置：2-1", A_COLOR),
    (2, "4-1", "4-1 工作台 BFF", "前置：2-1", A_COLOR),
    (3, "2-12", "2-12 批量导入", "前置：2-6＋2-10＋2-11", B_COLOR),
]
for layer, key, title, prereq, color in CARDS:
    dep_card(s, CX[key] - W / 2, LY[layer], W, H, title, prereq, color,
             fill=PAPER if color == MUTED else CREAM)

# 第4层合流卡（宽）+ 同步点标注放卡右侧（不压箭头落点）
box(s, Inches(2.9), Inches(LY[4]), Inches(8.0), Inches(0.8), fill=B_COLOR)
text(s, Inches(3.05), Inches(LY[4] + 0.04), Inches(7.7), Inches(0.72),
     [("4-20 账号运营三页 · 对照原型交付", 13, "FFFFFF", True),
      ("前置：2-10 · 2-11 · 2-12 · 4-1 · 2-5", 9, "FFFFFF", False)],
     anchor=MSO_ANCHOR.MIDDLE, spacing=1.08)
text(s, Inches(9.65), Inches(LY[3] + 0.22), Inches(1.9), Inches(0.28),
     "★同步点1", size=11, color=A_COLOR, bold=True)
text(s, Inches(11.02), Inches(LY[4] + 0.26), Inches(1.6), Inches(0.28),
     "★同步点2", size=11, color=A_COLOR, bold=True)

# 依赖箭头：落点分散到目标卡的不同边，避免同点汇合
BOT = lambda k, l: (CX[k], LY[l] + H)          # noqa: E731 卡底边中点
EDGES = [
    # (起点, 终点, 样式)  cp=关键路径粗黑线 dash=跨层虚线
    (BOT("1-4", 0), (CX["2-1"], LY[1]), "cp"),            # 1-4 → 2-1
    (BOT("1-6", 0), (CX["2-10"], LY[2]), "dash"),         # 1-6 → 2-10（竖直跨层）
    (BOT("2-1", 1), (CX["2-10"], LY[2]), "cp"),           # 2-1 → 2-10
    (BOT("2-1", 1), (CX["2-11"], LY[2]), ""),             # 2-1 → 2-11
    (BOT("2-1", 1), (CX["4-1"], LY[2]), ""),              # 2-1 → 4-1
    (BOT("2-6", 0), (CX["2-12"], LY[3]), "dash"),         # 2-6 → 2-12（竖直跨层，落顶边）
    (BOT("2-10", 2), (CX["2-12"] - W / 2, LY[3] + 0.20), ""),   # 2-10 → 2-12 左边沿
    (BOT("2-11", 2), (CX["2-12"] - W / 2, LY[3] + 0.50), ""),   # 2-11 → 2-12 左边沿低位
    (BOT("2-12", 3), (CX["2-12"], LY[4]), "cp"),          # 2-12 → 4-20 顶边
    (BOT("4-1", 2), (CX["4-1"], LY[4]), ""),              # 4-1 → 4-20 顶边（竖直）
    (BOT("2-5", 0), (10.5, LY[4]), "dash"),               # 2-5 → 4-20 右端（长跨层）
]
for (x0, y0), (x1, y1), style in EDGES:
    if style == "cp":
        kw = dict(color=CP_COLOR, width=2.25)
    else:
        kw = dict(color=INK_SOFT, width=1.0)
    arrow(s, x0, y0, x1, y1, dashed=(style == "dash"), **kw)

# ══ 第 2 页：盯三件事（保留）════════════════════════════════════════
s = add_slide(prs)
page_chrome(s, 2, "elephant-e-ai · 2a 冲刺排期")
text(s, Inches(0.55), Inches(0.66), Inches(12.2), Inches(0.5),
     "盯在这三件事，其余都可排后", size=30, color=INK, bold=True)
text(s, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.3),
     "两线各自动，见面只在两个同步点 · 9 月中 2a 硬提交为锚点", size=11, color=MUTED)

CW2, CH2 = 5.9, 2.35
cards = [
    (0.75, 1.75, "① 关键路径在 A 线", [
        ("1-4 → 2-1 → 2-10 → 2-12 → 4-20", 13, "274CA0", True),
        ("共 5 步最长，队友这条线不能窝工；", 11, INK_SOFT, False),
        ("B 线的 2-6 时间宽裕，空窗正好消化弹性件。", 11, INK_SOFT, False),
    ]),
    (6.95, 1.75, "② 两个硬同步点", [
        ("★1 · 2-12 开工前：B 的 2-6 + A 的 2-10 / 2-11 三方汇齐；", 11, INK_SOFT, False),
        ("★2 · 4-20 开工前：上图全部 + 2-5 审计待办落地。", 11, INK_SOFT, False),
        ("其余时间两线互不等待，从第一天起真并行。", 11, INK_SOFT, False),
    ]),
    (0.75, 4.35, "③ 风险与预案", [
        ("2-1 拖期 → B 线弹性件做完后接手 2-11 花名册", 11, INK_SOFT, False),
        ("（只依赖 2-1，不绑定 A 线本人）；", 11, INK_SOFT, False),
        ("1-6 必须赶在 2-10 之前，别让它卡进关键路径。", 11, INK_SOFT, False),
    ]),
    (6.95, 4.35, "④ 明确暂缓（防蔓延）", [
        ("2-2 MCP · 2-3 企微发群 · 2-4 成绩模板 · R3 布置全套", 11, INK_SOFT, False),
        ("等 9 月中 2a 硬提交完成后再排期。", 11, INK_SOFT, False),
    ]),
]
for x, y, title, lines in cards:
    box(s, Inches(x), Inches(y), Inches(CW2), Inches(CH2), fill=CREAM, line=LINE)
    text(s, Inches(x + 0.18), Inches(y + 0.16), Inches(CW2 - 0.36), Inches(0.35),
         title, size=19, color=INK, bold=True)
    text(s, Inches(x + 0.18), Inches(y + 0.62), Inches(CW2 - 0.36), Inches(CH2 - 0.78),
         lines, spacing=1.25)

target = os.path.normpath(OUT)
try:
    prs.save(target)
except PermissionError:  # 文件被 PowerPoint 占用时落 v2 名，不打断用户
    target = target.replace(".pptx", "-v2.pptx")
    prs.save(target)
print("OK saved:", target)
