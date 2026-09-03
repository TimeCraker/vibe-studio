# -*- coding: utf-8 -*-
"""
DDD (领域驱动设计) 架构知识点精要与落地实战 Deck 生成脚本
- 遵循 skills/ppt 规范：16:9 网格、7+1 封闭技术色板、严禁占位符
- 包含：封面、痛点对比、章节页、战略设计支柱、上下文映射架构图、战术设计元模型表格、
        聚合协作流向图、经典四层与依赖倒置架构图、贫血 vs 充血模型对比、避坑指南、总结金句、尾页
- 动画动效：COM 写入原生 PowerPoint 入场动效、级联微调
- 演讲者口播稿：每页内置 notes()
"""
import os
import sys

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(ROOT, "skills", "ppt", "templates"))

import primitives as P
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from animate import Anim

OUT_DIR = os.path.join(ROOT, "output")
os.makedirs(OUT_DIR, exist_ok=True)
OUT_PPTX = os.path.join(OUT_DIR, "ddd-architecture-deck.pptx")
OUT_PDF = os.path.join(OUT_DIR, "ddd-architecture-deck.pdf")

# ── 设计系统：架构师技术色板（冷蓝 / 藏青 / 琥珀点缀）───────────────────
P.THEMES["ddd_tech"] = dict(
    CORAL="2563EB",       # 核心高亮蓝 (Blue-600)
    CORAL_DEEP="1E40AF",  # 藏青蓝强调 (Blue-800)
    CREAM="F8FAFC",       # 极浅冷灰底 (Slate-50)
    INK="0F172A",         # 墨黑正文 (Slate-900)
    INK_SOFT="334155",    # 墨灰次级 (Slate-700)
    MUTED="64748B",       # 注释与页码灰 (Slate-500)
    LINE="E2E8F0",        # 分隔线与描边 (Slate-200)
)
P.use_theme("ddd_tech")
P.PAPER = "FFFFFF"

prs = P.new_deck()
anim = Anim(prs)

META = {
    "title": "DDD 架构深度拆解与实战精要",
    "subtitle": "从战略分界到战术落地的现代化软件设计范式",
    "kicker": "Domain-Driven Design Architecture",
    "tagline": "拆解复杂业务核心 · 守住领域不变性 · 依赖倒置架构实战",
    "domain": "architecture.ddd.guide",
    "version": "v1.0 Professional",
    "date": "2026-09",
    "author": "TimeCraker",
}

# ── P1: 封面 ──────────────────────────────────────────────────────────
s1 = P.slide_cover(prs, META)
P.notes(s1, "各位架构师与开发者好。今天我们用一套极其清晰、系统且不讲虚无概念的框架，"
            "彻底吃透领域驱动设计（DDD）。从战略设计的上下文划分，到战术设计的聚合根与代码心智，"
            "再到依赖倒置的分层落地，把 DDD 从天书变成工程生产力。")

# ── P2: 痛点与 Why DDD ───────────────────────────────────────────────
s2 = P.slide_versus(
    prs, 2, "WHY DDD", "为什么需要 DDD：打破贫血模型与大泥球架构的恶性循环",
    left=("传统数据驱动 / 三层贫血 CRUD 困境", [
        "以数据库表为中心：先画 ER 图再写代码，业务规则散落于各层",
        "Service 膨胀成上帝类：动辄数千行，充满 if-else 校验与逻辑外泄",
        "对象沦为数据传输袋：只有 getter/setter，无法保护业务规则的不变性",
        "牵一发而动全身：数据库字段变更引发全链路重构，无法进行独立演进",
    ]),
    right=("领域驱动设计 (DDD) 核心解法", [
        "以领域模型为中心：先挖掘业务不变性与概念边界，技术实现向后靠",
        "充血模型行为内聚：数据与行为高度内聚于实体，方法暴露业务意图",
        "依赖倒置隔绝基础设施：领域层是核心资产，零外部技术框架污染",
        "战略划定限界上下文：微服务与模块拆分有据可依，消除团队协作壁垒",
    ]),
    badge="Problem & Solution"
)
P.notes(s2, "我们先直面传统开发的核心痛点。大家常见的传统 MVC 三层架构，本质往往是‘贫血数据驱动’。"
            "开发者习惯先设计数据库表，再生成实体和 Service。结果 Service 沦为几千行的面条代码，"
            "实体只是纯数据的 Getter/Setter。而 DDD 的核心思想是颠覆这一现状，让领域模型成为第一公民，"
            "业务规则高度内聚于聚合内，技术细节反向依赖业务核心。")

# ── P3: 章节一：战略设计 ──────────────────────────────────────────────
s3 = P.slide_section(
    prs, 1, "战略设计：业务版图的解构与边界划定",
    points=[
        "统一语言（Ubiquitous Language）：消除研发与业务沟通的翻译损耗",
        "限界上下文（Bounded Context）：复杂系统分而治之的防御边界",
        "上下文映射（Context Map）：微服务拓扑与上下游协作协议",
    ],
    total=3,
    domain=META["domain"]
)
P.notes(s3, "第一部分：战略设计。战略设计解决的是‘方向与边界’的问题。"
            "如果战略边界画错了，战术设计代码写得再优雅，最终也只会变成一个巨大的分布式大泥球。")

# ── P4: 战略设计三大核心支柱 ──────────────────────────────────────────
s4 = P.slide_cards(
    prs, 4, "Strategic Pillars", "战略设计三大支柱：从概念对齐到边界收敛",
    [
        ("统一语言\nUbiquitous Language",
         "消除认知鸿沟：业务专家与开发必须共用同一套词汇体系\n"
         "代码即文档：类名、方法名直接对应业务术语，无黑话翻译\n"
         "边界内唯一：同一个词在不同上下文中可有不同专有含义"),
        ("限界上下文\nBounded Context",
         "明确语义边界：每个模型只在其所属上下文内具备唯一明确语义\n"
         "天然微服务边界：限界上下文是划分微服务物理边界的最佳标准\n"
         "团队自治：不同团队独立负责上下文，模型自治互不干扰"),
        ("上下文映射\nContext Mapping",
         "定义协作关系：上下游依赖（U/D）、客户/供应者（C/S）\n"
         "防腐层 (ACL)：隔离外部脏模型，保护本上下文核心领域纯洁\n"
         "共享内核与开放主机：精细化控制系统间集成成本与演进自由度"),
    ],
    cols=3,
    card_h=3.8
)
P.notes(s4, "战略设计有三大支柱。第一是统一语言，严禁用研发的思维去教业务，或者开发自己造一套翻译字典；"
            "第二是限界上下文，解决的是‘多义性’问题，比如同一个‘用户’，在风控上下文是‘征信主体’，在营销上下文是‘流量受众’，"
            "绝不能硬捏成一张大表；第三是上下文映射，理清系统间的上下游权力关系和防腐机制。")

# ── P5: 限界上下文架构图（Context Map 矢量拓扑流程）───────────────────
def slide_context_map_diagram(prs, idx, label, title):
    s = P.add_slide(prs)
    P.page_chrome(s, idx, label, badge="Architecture Diagram")
    P.text(s, Inches(P.MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
           title, 30, P.INK, True).name = "title"

    # 顶部说明条
    P.box(s, Inches(P.MARGIN), Inches(1.85), Inches(12.23), Inches(0.48), fill=P.CREAM, line=P.LINE).name = "ctx_desc_box"
    P.text(s, Inches(P.MARGIN + 0.2), Inches(1.92), Inches(11.8), Inches(0.35),
           "电商微服务协同拓扑：订单上下文（核心域）通过 防腐层(ACL) 隔离商品、支付与履约外部模型",
           11, P.INK_SOFT, True, font=P.FONT_CN).name = "ctx_desc_t"

    # 上游1：商品目录上下文 (Catalog Context)
    c1_x, c1_y, cw, ch = 0.55, 2.6, 3.2, 1.8
    P.box(s, Inches(c1_x), Inches(c1_y), Inches(cw), Inches(ch), fill=P.PAPER, line=P.LINE).name = "ctx1box"
    P.box(s, Inches(c1_x), Inches(c1_y), Inches(cw), Inches(0.08), fill="64748B").name = "ctx1bar"
    P.text(s, Inches(c1_x + 0.2), Inches(c1_y + 0.2), Inches(cw - 0.4), Inches(0.35),
           "商品上下文 (Upstream)", 13, P.INK, True).name = "ctx1t"
    P.text(s, Inches(c1_x + 0.2), Inches(c1_y + 0.6), Inches(cw - 0.4), Inches(1.1),
           ["• 关注：SPU/SKU规格、类目、价格", "• 角色：上游生产者 (U)", "• 接口：OHS/PL 开放标准 API"],
           10, P.INK_SOFT, spacing=1.3).name = "ctx1d"

    # 核心：订单上下文 (Order Context) - 重点高亮
    c2_x, c2_y, c2w, c2h = 4.35, 2.6, 4.6, 3.8
    P.box(s, Inches(c2_x), Inches(c2_y), Inches(c2w), Inches(c2h), fill=P.CREAM, line=P.CORAL).name = "ctx2box"
    P.box(s, Inches(c2_x), Inches(c2_y), Inches(c2w), Inches(0.12), fill=P.CORAL).name = "ctx2bar"
    P.text(s, Inches(c2_x + 0.25), Inches(c2_y + 0.25), Inches(c2w - 0.5), Inches(0.35),
           "订单上下文 (Core Domain · 核心域)", 15, P.CORAL_DEEP, True).name = "ctx2t"

    # 订单内部的防腐层 (ACL)
    acl_w, acl_h = c2w - 0.5, 0.95
    P.box(s, Inches(c2_x + 0.25), Inches(c2_y + 0.75), Inches(acl_w), Inches(acl_h), fill=P.PAPER, line="3B82F6").name = "acl1box"
    P.text(s, Inches(c2_x + 0.35), Inches(c2_y + 0.85), Inches(acl_w - 0.2), Inches(0.3),
           "防腐层 Anti-Corruption Layer (ACL)", 11, "1D4ED8", True, font=P.FONT_MONO).name = "acl1t"
    P.text(s, Inches(c2_x + 0.35), Inches(c2_y + 1.15), Inches(acl_w - 0.2), Inches(0.5),
           "将外部 SKU/促销模型翻译为内部 OrderItem 快照，阻断外部变动污染核心领域", 9.5, P.INK_SOFT).name = "acl1d"

    # 订单核心聚合
    P.box(s, Inches(c2_x + 0.25), Inches(c2_y + 1.9), Inches(acl_w), Inches(1.6), fill=P.PAPER, line=P.LINE).name = "ord_core_box"
    P.text(s, Inches(c2_x + 0.35), Inches(c2_y + 2.05), Inches(acl_w - 0.2), Inches(0.3),
           "订单聚合根 (Order Aggregate Root)", 12, P.INK, True).name = "ord_core_t"
    P.text(s, Inches(c2_x + 0.35), Inches(c2_y + 2.4), Inches(acl_w - 0.2), Inches(1.0),
           ["• 封装状态机：待支付、已确认、已关闭", "• 守护不变性：订单总金额与明细校验", "• 发布领域事件：OrderCreated, OrderPaid"],
           10, P.INK_SOFT, spacing=1.3).name = "ord_core_d"

    # 上游/下游：支付上下文 (Payment Context)
    c3_x, c3_y = 9.5, 2.6
    P.box(s, Inches(c3_x), Inches(c3_y), Inches(cw), Inches(ch), fill=P.PAPER, line=P.LINE).name = "card2box"
    P.box(s, Inches(c3_x), Inches(c3_y), Inches(cw), Inches(0.08), fill="047857").name = "card2bar"
    P.text(s, Inches(c3_x + 0.2), Inches(c3_y + 0.2), Inches(cw - 0.4), Inches(0.35),
           "支付上下文 (Downstream/Peer)", 13, P.INK, True).name = "card2t"
    P.text(s, Inches(c3_x + 0.2), Inches(c3_y + 0.6), Inches(cw - 0.4), Inches(1.1),
           ["• 关注：收单凭证、交易渠道、退款", "• 协议：基于领域事件异步解耦", "• 订阅 OrderCreated 事件创建账单"],
           10, P.INK_SOFT, spacing=1.3).name = "card2d"

    # 下游：仓储履约上下文 (Fulfillment Context)
    c4_x, c4_y = 9.5, 4.6
    P.box(s, Inches(c4_x), Inches(c4_y), Inches(cw), Inches(ch), fill=P.PAPER, line=P.LINE).name = "card3box"
    P.box(s, Inches(c4_x), Inches(c4_y), Inches(cw), Inches(0.08), fill="B45309").name = "card3bar"
    P.text(s, Inches(c4_x + 0.2), Inches(c4_y + 0.2), Inches(cw - 0.4), Inches(0.35),
           "履约仓配上下文 (Fulfillment)", 13, P.INK, True).name = "card3t"
    P.text(s, Inches(c4_x + 0.2), Inches(c4_y + 0.6), Inches(cw - 0.4), Inches(1.1),
           ["• 关注：拣货单、运单路由、出库", "• 角色：下游跟随者 (D)", "• 响应 OrderPaid 触发仓储打包"],
           10, P.INK_SOFT, spacing=1.3).name = "card3d"

    # 左右连接线/箭头文字（使用高对比度颜色）
    P.text(s, Inches(3.75), Inches(3.2), Inches(0.6), Inches(0.4), "→", 24, P.CORAL, True, PP_ALIGN.CENTER)
    P.text(s, Inches(8.95), Inches(3.2), Inches(0.55), Inches(0.4), "→", 24, "047857", True, PP_ALIGN.CENTER)
    P.text(s, Inches(8.95), Inches(5.2), Inches(0.55), Inches(0.4), "→", 24, "B45309", True, PP_ALIGN.CENTER)

    # 底部说明
    P.text(s, Inches(P.MARGIN), Inches(6.6), Inches(12.2), Inches(0.4),
           "架构启示：限界上下文之间禁止直接数据库跨库联表查询；必须通过 RPC/API 结合防腐层适配，或通过领域事件实现最终一致性。",
           10, P.MUTED, True).name = "ctx_foot"

    return s

s5 = slide_context_map_diagram(prs, 5, "Context Map", "限界上下文协作与防腐层 (ACL) 架构全景")
P.notes(s5, "大家看第 5 页的架构全景图。在微服务架构下，订单是核心域，而商品、支付、履约是外部或下游上下文。"
            "重点看订单内部的防腐层（ACL）：很多团队犯的错误是把商品服务的 DTO 直接当成订单实体使用，"
            "导致商品微服务一旦加字段改结构，订单全崩盘。通过 ACL 防腐层转换成内部领域模型，"
            "实现了真正的高内聚与自主演进。")

# ── P6: 章节二：战术设计 ──────────────────────────────────────────────
s6 = P.slide_section(
    prs, 2, "战术设计：领域模型的核心构件与心智模式",
    points=[
        "实体 vs 值对象：唯一标识与无状态不可变的取舍",
        "聚合与聚合根：高内聚的不变性边界与事务锁单元",
        "领域服务与领域事件：无状态逻辑承载与异步解耦",
        "仓储与工厂：领域生命周期的装配与持久化倒置",
    ],
    total=3,
    domain=META["domain"]
)
P.notes(s6, "第二部分：战术设计。战术设计提供了丰富的领域模式构建块，"
            "用来在代码层面落地面向对象设计的精髓。我们将重点搞清楚实体、值对象和聚合根的本质区别。")

# ── P7: 战术核心元模型真表格 ──────────────────────────────────────────
s7 = P.slide_table(
    prs, 7, "Building Blocks", "战术设计核心元模型矩阵：职责与代码特征",
    cols=["元模型概念", "核心定义与职责", "标识与状态特性", "典型应用场景", "代码实现黄金法则"],
    rows=[
        ["实体 (Entity)", "具有唯一业务标识与连续生命周期的对象", "依赖 ID 区分平等性，属性随业务变化", "订单 (Order)、用户 (User)", "以业务行为命名方法，禁止无意义 Setter"],
        ["值对象 (Value Object)", "用于描述事物特征且无概念标识的属性度量", "完全不可变 (Immutable)，基于属性判等", "地址 (Address)、金额 (Money)", "属性只读，任何修改均返回新值对象实例"],
        ["聚合根 (Aggregate Root)", "聚合的入口与主脑，维护聚合内部业务一致性", "外部访问聚合内部实体的唯一网关", "订单聚合维护 OrderItem 与收件信息", "聚合间只引用 ID 不引用对象，单事务只改一个聚合"],
        ["领域服务 (Domain Service)", "跨越多个聚合、不属于任何单一实体的业务逻辑", "纯行为、无状态的业务操作接口", "转账服务 (A扣减+B增加)、复合折扣计算", "避免滥用沦为贫血 Service，先看是否能放聚合根"],
        ["领域事件 (Domain Event)", "记录领域内已经发生并对业务有重要意义的事实", "只读、携带发生时间与关键事实载荷", "OrderPaidEvent, UserRegisteredEvent", "强一致转最终一致性，通过 MQ 驱动下游系统解耦"],
        ["仓储 (Repository)", "提供对聚合根集合形态的访问抽象，隔离持久化", "领域层定义接口，基础设施层负责实现", "OrderRepository.save(Order)", "只针对聚合根设计仓储，一个聚合对应一个仓储"],
    ],
    col_ws=[0.16, 0.28, 0.22, 0.16, 0.18],
    row_h=0.68
)
P.notes(s7, "这是整个战术设计最核心的一张对照表。请大家牢记：实体看唯一标识（哪怕属性都一样，ID 不同也是两个人）；"
            "值对象看属性（两张百元大钞只要面值相同，就没有区别，值对象必须设计成不可变对象）；"
            "聚合根则是整个聚合的防线，外部任何调用禁止绕过聚合根直接修改内部明细！")

# ── P8: 战术核心交互链路（流程图）────────────────────────────────────
s8 = P.slide_chain(
    prs, 8, "TACTICAL FLOW", "战术设计运行流：从外部指令到聚合状态变更的完整链路",
    nodes=[
        ("1. 用户指令", "Command/DTO 请求"),
        ("2. 应用服务", "AppService 事务编排"),
        ("3. 仓储加载", "Repo 还原聚合根"),
        ("4. 聚合根执行", "充血方法守护不变性"),
        ("5. 仓储保存", "持久化聚合状态"),
        ("6. 发布事件", "Domain Event 广播"),
    ],
    hi=3,
    sub_title="关键组件分工细则",
    subs=[
        ("应用层编排", "AppService 负责鉴权、开启事务、调用仓储加载聚合，绝不写核心业务分支计算"),
        ("聚合根守护核心", "Order.pay() 内部校验状态是否待支付，扣减优惠，标记支付完成并记录事件"),
        ("事件驱动下游", "通过 Outbox 模式可靠落库并投递 MQ，通知发货、开票，实现系统解耦与高吞吐"),
    ]
)
P.notes(s8, "我们用一条清晰的生命周期流水线理顺各组件的关系：Controller 接收到请求转化为 Command；"
            "应用服务开启事务，调用仓储将订单聚合根完整加载进内存；第四步是核心——调用聚合根的充血业务方法，"
            "内部执行状态转换和业务校验；然后保存回仓储；最后由领域事件驱动下游。应用服务做编排，聚合根做决策！")

# ── P9: 章节三：分层架构与依赖倒置 ──────────────────────────────────
s9 = P.slide_section(
    prs, 3, "分层架构与依赖倒置：保护领域内核的技术屏障",
    points=[
        "经典 DDD 四层架构：接口层、应用层、领域层、基础设施层",
        "依赖倒置原则 (DIP)：领域层零技术依赖的秘密",
        "六边形架构与洋葱模型：端口与适配器的工程具象化",
    ],
    total=3,
    domain=META["domain"]
)
P.notes(s9, "第三部分：架构分层。为什么很多人写着写着，领域层还是引入了 Spring 注解、MyBatis 注解和各种外部 SDK？"
            "因为他们没有真正理解依赖倒置原则。接下来我们详细拆解经典的四层架构与 DIP 机制。")

# ── P10: DDD 经典四层与依赖倒置架构图（专业矢量绘制）─────────────────
def slide_layered_architecture_diagram(prs, idx, label, title):
    s = P.add_slide(prs)
    P.page_chrome(s, idx, label, badge="DIP Architecture")
    P.text(s, Inches(P.MARGIN), Inches(0.95), Inches(12.2), Inches(0.7),
           title, 30, P.INK, True).name = "title"

    # 左侧：分层结构主视区 (宽 7.6)
    lx, ly, lw = 0.55, 1.85, 7.6

    # 第 1 层：用户接口层
    P.box(s, Inches(lx), Inches(ly), Inches(lw), Inches(0.95), fill=P.PAPER, line=P.LINE).name = "l1box"
    P.box(s, Inches(lx), Inches(ly), Inches(0.12), Inches(0.95), fill="64748B").name = "l1bar"
    P.text(s, Inches(lx + 0.3), Inches(ly + 0.12), Inches(3.0), Inches(0.3),
           "用户接口层 (Interfaces / UI)", 12, P.INK, True).name = "l1t"
    P.text(s, Inches(lx + 0.3), Inches(ly + 0.45), Inches(lw - 0.5), Inches(0.4),
           "HTTP RESTful API, RPC Facade, GraphQL, MQ Consumer, WebSocket 协议适配", 10, P.INK_SOFT).name = "l1d"

    # 箭头 1
    P.text(s, Inches(lx + lw / 2 - 0.2), Inches(ly + 0.95), Inches(0.4), Inches(0.3),
           "↓", 14, P.CORAL, True, PP_ALIGN.CENTER)

    # 第 2 层：应用层
    ly2 = ly + 1.25
    P.box(s, Inches(lx), Inches(ly2), Inches(lw), Inches(1.05), fill=P.PAPER, line=P.LINE).name = "l2box"
    P.box(s, Inches(lx), Inches(ly2), Inches(0.12), Inches(1.05), fill="3B82F6").name = "l2bar"
    P.text(s, Inches(lx + 0.3), Inches(ly2 + 0.12), Inches(3.5), Inches(0.3),
           "应用层 (Application Layer)", 12, "1D4ED8", True).name = "l2t"
    P.text(s, Inches(lx + 0.3), Inches(ly2 + 0.45), Inches(lw - 0.5), Inches(0.55),
           "用例编排 (Use Case Orchestration)、事务控制 (@Transactional)、安全鉴权、DTO 转换", 10, P.INK_SOFT).name = "l2d"

    # 箭头 2
    P.text(s, Inches(lx + lw / 2 - 0.2), Inches(ly2 + 1.05), Inches(0.4), Inches(0.3),
           "↓", 14, P.CORAL, True, PP_ALIGN.CENTER)

    # 第 3 层：领域层（核心，用 CREAM + 蓝色粗边高亮）
    ly3 = ly2 + 1.35
    P.box(s, Inches(lx), Inches(ly3), Inches(lw), Inches(1.35), fill=P.CREAM, line=P.CORAL).name = "l3box"
    P.box(s, Inches(lx), Inches(ly3), Inches(0.18), Inches(1.35), fill=P.CORAL).name = "l3bar"
    P.text(s, Inches(lx + 0.3), Inches(ly3 + 0.15), Inches(5.5), Inches(0.3),
           "领域层 (Domain Layer) · 系统的绝对核心", 13, P.CORAL_DEEP, True).name = "l3t"
    P.text(s, Inches(lx + 0.3), Inches(ly3 + 0.52), Inches(lw - 0.5), Inches(0.75),
           ["• 实体、值对象、聚合根：纯内聚业务状态机与不变规则",
            "• 仓储接口 (OrderRepository)：定义按 ID 存取契约（无 SQL 痕迹）",
            "• 严格规则：纯 POJO/内存模型，零外部框架依赖，不依赖基础设施"],
           10, P.INK, spacing=1.2).name = "l3d"

    # 反转箭头（依赖倒置核心）
    P.text(s, Inches(lx + lw / 2 - 0.2), Inches(ly3 + 1.35), Inches(0.4), Inches(0.3),
           "↑", 16, "047857", True, PP_ALIGN.CENTER)
    P.text(s, Inches(lx + lw / 2 + 0.3), Inches(ly3 + 1.38), Inches(3.0), Inches(0.3),
           "依赖倒置 (DIP): 实现依赖接口", 9.5, "047857", True, font=P.FONT_MONO)

    # 第 4 层：基础设施层
    ly4 = ly3 + 1.65
    P.box(s, Inches(lx), Inches(ly4), Inches(lw), Inches(1.05), fill=P.PAPER, line=P.LINE).name = "l4box"
    P.box(s, Inches(lx), Inches(ly4), Inches(0.12), Inches(1.05), fill="047857").name = "l4bar"
    P.text(s, Inches(lx + 0.3), Inches(ly4 + 0.12), Inches(4.5), Inches(0.3),
           "基础设施层 (Infrastructure Layer)", 12, "047857", True).name = "l4t"
    P.text(s, Inches(lx + 0.3), Inches(ly4 + 0.45), Inches(lw - 0.5), Inches(0.55),
           "仓储具体实现 (Spring Data / MyBatis / Redis)、MQ 生产者、第三方服务 SDK 适配器", 10, P.INK_SOFT).name = "l4d"

    # 右侧：依赖倒置原理解析卡片 (宽 4.3)
    rx, ry, rw = 8.45, 1.85, 4.33
    P.box(s, Inches(rx), Inches(ry), Inches(rw), Inches(4.8), fill=P.CREAM, line=P.LINE).name = "dip_box"
    P.box(s, Inches(rx), Inches(ry), Inches(rw), Inches(0.08), fill=P.CORAL).name = "dip_bar"
    P.text(s, Inches(rx + 0.25), Inches(ry + 0.2), Inches(rw - 0.5), Inches(0.35),
           "DIP 依赖倒置核心心智", 14, P.CORAL_DEEP, True).name = "dip_t"

    principles = [
        ("高层不依赖低层", "传统三层架构是：应用 -> 业务 -> 数据库，业务层直接依赖数据库驱动。而在 DDD 中，两层都依赖抽象。"),
        ("抽象不依赖细节", "领域层定义 OrderRepository 接口，完全不知晓底层是 MySQL、MongoDB 还是内存缓存。"),
        ("可测试性爆发", "由于领域层没有任何 Spring 上下文或数据库连接，单测聚合根只需毫秒级的纯 JVM 单元测试，效率提升 100 倍。"),
        ("灵活替换底座", "将存储从 MySQL 切到 TiDB 或 Redis，只需在基础设施层写新的适配器，领域层代码一行不改。"),
    ]
    cur_y = ry + 0.65
    for h, d in principles:
        P.box(s, Inches(rx + 0.25), Inches(cur_y + 0.06), Inches(0.12), Inches(0.12), fill=P.CORAL).name = "dip_dot"
        P.text(s, Inches(rx + 0.45), Inches(cur_y), Inches(rw - 0.7), Inches(0.28), h, 11, P.INK, True)
        P.text(s, Inches(rx + 0.45), Inches(cur_y + 0.28), Inches(rw - 0.7), Inches(0.65), d, 9.5, P.INK_SOFT, spacing=1.2)
        cur_y += 0.98

    return s

s10 = slide_layered_architecture_diagram(prs, 10, "Layers & DIP", "经典四层架构与依赖倒置原则 (DIP) 落地拆解")
P.notes(s10, "请看第 10 页。这是 DDD 架构最漂亮的地方：注意下半部分的绿色箭头！"
             "在传统开发中，业务层依赖数据访问层；但在 DDD 中，仓储接口声明在领域层（高层），"
             "MyBatis 或 JPA 的具体实现类写在基础设施层（低层）。这彻底切断了数据库技术对业务模型的侵入，"
             "让业务逻辑可以脱离任何外部中间件独立存活和极速测试！")

# ── P11: 贫血模型 VS 充血模型深度剖析 ─────────────────────────────────
s11 = P.slide_versus(
    prs, 11, "MODEL COMPARISON", "充血模型 VS 贫血模型：代码组织与维护成本质变",
    left=("贫血模型 (Anemic Domain Model)", [
        "类结构：只有属性字段 + getter/setter，实质只是个 C 风格的 struct",
        "逻辑外泄：校验、计算、状态流转全部写在 Service 的冗长方法里",
        "无规则保护：外部任何代码都可直接调用 .setStatus('PAID') 随意篡改状态",
        "复用性极差：不同开发在不同接口中重复抄写相同的业务校验逻辑",
        "维护噩梦：新同学不敢改 Service，只能在末尾继续叠加 if-else 补丁",
    ]),
    right=("充血模型 (Rich Domain Model)", [
        "类结构：数据与业务行为高度内聚，私有字段，对外仅暴露业务意图方法",
        "行为自洽：order.cancel(reason) 内部自动完成状态前置校验与事件挂载",
        "保护不变性：外部无法绕过规则直接修改数据，聚合永远处于合法状态",
        "代码即业务：读聚合代码如同阅读产品 PRD 规则手册，认知成本极低",
        "高可靠演化：业务规则变更只需修改聚合内一处逻辑，自动全量生效",
    ]),
    badge="Deep Dive"
)
P.notes(s11, "很多团队说自己上了 DDD，一打开代码还是纯贫血模型。贫血模型最大的危害在于‘无法保护业务不变性’。"
            "任何人只要拿到了 Order 对象，都可以调用 setStatus 把它改成已支付，完全不检查订单当前是否超时。而在充血模型中，"
            "想要改变状态必须调用 order.pay(paymentInfo)，方法内部自己守门，守不住就抛出领域异常。")

# ── P12: 落地避坑指南与反模式警示 ─────────────────────────────────────
s12 = P.slide_cards(
    prs, 12, "Pitfalls & Patterns", "DDD 落地避坑指南：架构师常踩的四大雷区",
    [
        ("雷区一：大聚合反模式",
         "现象：把所有关联对象塞进一个超级聚合\n"
         "危害：并发悲观锁导致系统坍塌与 OOM\n"
         "药方：小聚合设计，跨聚合最终一致"),
        ("雷区二：杀鸡用牛刀",
         "现象：简单字典或单表报表也强套四层\n"
         "危害：产生海量无用转换代码，研发低效\n"
         "药方：按业务复杂度分流，CRUD 走快道"),
        ("雷区三：形式主义伪 DDD",
         "现象：仅建 domain 包，里面全为贫血 Getter\n"
         "危害：未享高内聚优势，反增包层级成本\n"
         "药方：建立代码规范，严禁 Service 写判定"),
        ("雷区四：跨聚合事务蔓延",
         "现象：一个事务同时修改订单、库存多个聚合\n"
         "危害：锁竞争激烈，微服务拆分难以剥离\n"
         "药方：严格单事务单聚合，领域事件异步解耦"),
    ],
    cols=4,
    card_h=3.8
)
P.notes(s12, "第 12 页我们总结了落地 DDD 的四条高压红线。第一不要设计超级大聚合，很多新人把‘整个订单+商品+店铺+用户’做成一个聚合，"
            "直接把数据库锁死；第二区分业务场景，纯报表查询走 CQRS 的读通道，别硬套 DDD；第三要谨防形式主义的伪 DDD；"
            "第四牢记：一个事务只修改一个聚合根！")

# ── P13: 架构金句页 ──────────────────────────────────────────────────
s13 = P.slide_quote(
    prs,
    "软件开发的核心复杂度，不在于所使用的编程语言或技术框架，\n而在于问题域本身的业务逻辑与规则。",
    "Eric Evans · 《领域驱动设计：软件核心复杂性应对之道》"
)
P.notes(s13, "最后分享 Eric Evans 在 DDD 奠基名作中的这句至理名言。框架每五年换一代，而业务领域的内核知识沉淀，"
            "才是软件工程中最有价值的数字资产。")

# ── P14: 尾页 ────────────────────────────────────────────────────────
s14 = P.slide_closing(
    prs, META,
    footer_lines=[
        "演讲主题：领域驱动设计 (DDD) 架构精要与工程落地",
        "讲师：TimeCraker · AI 时代独立架构师",
        "交流社区与资源归档：architecture.ddd.guide",
    ],
    slogan="让代码真实映射业务知识，构筑高内聚的领域内核。"
)
P.notes(s14, "今天的分享就到这里，感谢大家的聆听！欢迎交流与提问。")

# ── 动画编排与保存 ───────────────────────────────────────────────────
print("[INFO] Applying automatic animation sequence...")
anim.auto_deck()
P.set_transition(prs, "fade")
P.auto_show(prs)

prs.save(OUT_PPTX)
print(f"[OK] Saved pptx to: {OUT_PPTX}")

# 写入 COM 原生 PowerPoint 动画
try:
    anim.apply(OUT_PPTX)
    print(f"[OK] Applied COM animations successfully.")
except Exception as e:
    print(f"[WARN] COM animation application skipped or failed: {e}")

print("[DONE] Deck generation completed.")
